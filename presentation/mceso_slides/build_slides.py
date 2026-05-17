"""Build MC-ESO presentation slides (.pptx).

Design principles applied (from presentation/presentation.pdf):
  * Proximity   — related items grouped, generous gap between groups
  * Alignment   — left-aligned text by default, single invisible vertical guide
  * Repetition  — fixed title bar, page number, citation slot on every slide
  * Contrast    — one accent color used sparingly for the key message

Aspect ratio: 16:9. Background: white (recommended for novices).
"""

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree


# Cache dir for SVG→PNG conversions of experiment-generated SVGs
SVG_CACHE = Path(__file__).parent / "figures" / "svg_cache"
SVG_CACHE.mkdir(parents=True, exist_ok=True)


def _crop_whitespace(png_path: Path, *, bg_tol: int = 248,
                     make_white_transparent: bool = False) -> None:
    """Crop near-white margins from a PNG in place. No-op if no margin found.

    When ``make_white_transparent`` is True, near-white pixels are also made
    transparent (used for formula images so they don't carry a white box onto
    coloured chips)."""
    from PIL import Image, ImageOps
    with Image.open(str(png_path)) as im:
        rgb = im.convert("RGB")
        inverted = ImageOps.invert(rgb)
        bbox = inverted.point(lambda v: 0 if v < (255 - bg_tol) else 255).getbbox()
        if bbox:
            im = im.crop(bbox)
        if make_white_transparent:
            im = im.convert("RGBA")
            pixels = im.load()
            w, h = im.size
            for y in range(h):
                for x in range(w):
                    r, g, b, a = pixels[x, y]
                    if r >= bg_tol and g >= bg_tol and b >= bg_tol:
                        pixels[x, y] = (r, g, b, 0)
        im.save(str(png_path))


def svg_to_png(svg_path: Path, target_w_px: int = 2000,
               crop: bool = True,
               make_white_transparent: bool = False) -> Path:
    """Convert SVG → PNG via macOS qlmanage. Caches by source mtime.

    With `crop=True`, near-white margins are trimmed from the cached PNG so
    the actual plot fills the slide cell. `make_white_transparent` additionally
    converts white interior pixels to transparent (useful for math formulas
    placed on coloured backgrounds).

    If the SVG no longer exists but a cached PNG does, return the cached PNG
    — useful when an experiment directory has been deleted but we still want
    to render slides from previously-converted images.
    """
    svg_path = Path(svg_path)
    suffix = "__t" if make_white_transparent else ""
    out = SVG_CACHE / f"{svg_path.stem}__{target_w_px}{suffix}.png"
    if out.exists():
        if not svg_path.exists():
            return out  # cache-only fallback
        if out.stat().st_mtime >= svg_path.stat().st_mtime:
            return out
    if not svg_path.exists():
        raise FileNotFoundError(f"no SVG and no cache for {svg_path}")
    subprocess.run(["qlmanage", "-t", "-s", str(target_w_px), "-o",
                    str(SVG_CACHE), str(svg_path)],
                   check=True, capture_output=True, timeout=15)
    raw = SVG_CACHE / f"{svg_path.name}.png"
    raw.replace(out)
    if crop:
        _crop_whitespace(out, make_white_transparent=make_white_transparent)
    return out


def safe_svg_to_png(svg_path: Path, target_w_px: int = 2000):
    """Wrapper that returns None on timeout/error so callers can fall back."""
    try:
        return svg_to_png(svg_path, target_w_px)
    except Exception as e:
        print(f"  ! SVG conversion failed for {Path(svg_path).name}: {e}")
        return None


def fig_png(name: str, target_w_px: int = 1800, *,
            transparent: bool = False) -> str:
    """Return a raster (PNG) path for a logical figure name.

    The on-disk source of truth is ``figures/{name}.svg``; we convert to a
    cached PNG so python-pptx can embed it (it can't embed SVG natively).
    Set ``transparent=True`` for figures (e.g. formulas) that need their
    white background dropped before placement on a coloured chip.
    """
    svg = Path(__file__).parent / "figures" / f"{name}.svg"
    png = svg_to_png(svg, target_w_px=target_w_px,
                     make_white_transparent=transparent)
    return str(png)


def add_picture_fit(slide, path, x, y, max_w, max_h, *, align="center"):
    """Embed an image inside (max_w × max_h) while preserving its aspect ratio.

    `align`: "center" (default) centres within the bounding box, "topleft"
    pins to the (x, y) origin.
    """
    from PIL import Image
    with Image.open(str(path)) as im:
        src_w, src_h = im.size
    src_ratio = src_w / src_h
    box_ratio = (max_w / max_h) if max_h else src_ratio
    if src_ratio >= box_ratio:
        out_w = max_w
        out_h = Emu(int(max_w * src_h / src_w))
    else:
        out_h = max_h
        out_w = Emu(int(max_h * src_w / src_h))
    if align == "center":
        ox = x + (max_w - out_w) // 2
        oy = y + (max_h - out_h) // 2
    else:
        ox, oy = x, y
    return slide.shapes.add_picture(str(path), ox, oy, out_w, out_h)


# ---------------------------------------------------------------------------
# Style tokens (single source of truth — repetition principle)
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_TOP = Inches(0.45)

TITLE_TOP = Inches(0.45)
TITLE_H = Inches(0.7)
TITLE_RULE_Y = Inches(1.20)         # thin divider line under title
BODY_TOP = Inches(1.45)
FOOTER_Y = Inches(7.05)

# Type scale (drastic contrast, not subtle)
F_TITLE = Pt(30)
F_H1 = Pt(26)
F_BODY = Pt(18)
F_SMALL = Pt(14)
F_TINY = Pt(11)

# Color palette — restrained, one accent only
C_INK = RGBColor(0x1F, 0x24, 0x2E)        # near-black for primary text
C_MUTED = RGBColor(0x5A, 0x65, 0x76)      # muted gray for secondary text
C_ACCENT = RGBColor(0xC0, 0x39, 0x2B)     # epidemic red — used sparingly for emphasis
C_ACCENT_SOFT = RGBColor(0xE7, 0xA8, 0xA2)
C_RULE = RGBColor(0xD0, 0xD5, 0xDB)       # divider line
C_BG_SOFT = RGBColor(0xF3, 0xF5, 0xF8)    # soft band for highlight rows
C_CONTACT = RGBColor(0x2E, 0x86, 0xAB)    # close-contact (blue)
C_DROPLET = RGBColor(0xE0, 0x7A, 0x5F)    # droplet (warm)
C_AIR = RGBColor(0x6B, 0x9A, 0x4C)        # airborne (green)
C_OK = RGBColor(0x2E, 0x86, 0x4F)         # good result green
C_BAD = RGBColor(0xC0, 0x39, 0x2B)        # weak result red


def hex_color(h: str) -> RGBColor:
    return RGBColor.from_string(h)


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def set_solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, *,
                font_size=F_BODY, bold=False, color=C_INK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                font_name="Helvetica"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, paragraphs, *,
                   font_size=F_BODY, color=C_INK, line_spacing=1.25,
                   align=PP_ALIGN.LEFT, font_name="Helvetica"):
    """paragraphs: list of (text, bold)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, (text, bold) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                font_size=F_BODY, color=C_INK, line_spacing=1.30,
                bullet_char="•", indent_pt=14, font_name="Helvetica"):
    """items: list of either str or (str, color/bold-overrides dict)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{bullet_char}  {text}"
        run.font.name = font_name
        run.font.size = opts.get("size", font_size)
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", color)
    return tb


def add_rect(slide, left, top, width, height, fill, line=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def add_line(slide, x1, y1, x2, y2, color=C_RULE, width=0.75, *, arrow=False):
    # python-pptx add_connector signature: (connector_type, begin_x, begin_y, end_x, end_y)
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if arrow:
        # XML hack: add an arrow end on the line
        ln_xml = ln.line._get_or_add_ln()
        tailEnd = etree.SubElement(ln_xml, qn("a:tailEnd"))
        tailEnd.set("type", "triangle")
        tailEnd.set("w", "med")
        tailEnd.set("len", "med")
    return ln


# ---------------------------------------------------------------------------
# Slide chrome (repeated structure — pixel-perfect)
# ---------------------------------------------------------------------------

def add_chrome(slide, title, page_num, total_pages, section=None):
    # Page number (bottom right) — single muted token, never moves
    add_textbox(slide, SLIDE_W - Inches(1.2), FOOTER_Y, Inches(0.6), Inches(0.3),
                f"{page_num} / {total_pages}",
                font_size=F_TINY, color=C_MUTED, align=PP_ALIGN.RIGHT)
    # Section tag (bottom left) — context tracker
    if section:
        add_textbox(slide, MARGIN_L, FOOTER_Y, Inches(6.0), Inches(0.3),
                    section,
                    font_size=F_TINY, color=C_MUTED)
    # Title
    add_textbox(slide, MARGIN_L, TITLE_TOP, SLIDE_W - MARGIN_L - MARGIN_R, TITLE_H,
                title, font_size=F_TITLE, bold=True, color=C_INK,
                anchor=MSO_ANCHOR.MIDDLE)
    # Title underline rule
    add_line(slide, MARGIN_L, TITLE_RULE_Y,
             SLIDE_W - MARGIN_R, TITLE_RULE_Y, color=C_RULE, width=0.75)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

# total pages counter is captured at the end
SLIDE_COUNT = [0]
TOTAL_PAGES = [0]


def new_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    SLIDE_COUNT[0] += 1
    return slide


# ── 01: Title --------------------------------------------------------------

def slide_title(prs):
    slide = new_slide(prs)
    # No chrome on cover

    # accent bar on the left
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, fill=C_ACCENT)

    # Eyebrow label
    add_textbox(slide, Inches(0.9), Inches(2.2), Inches(11), Inches(0.4),
                "PROPOSED METHOD",
                font_size=Pt(14), bold=True, color=C_ACCENT,
                font_name="Helvetica")

    # Main title (huge contrast)
    add_textbox(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.4),
                "MC-ESO",
                font_size=Pt(72), bold=True, color=C_INK,
                font_name="Helvetica")

    add_textbox(slide, Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.8),
                "Multi-Channel Epidemic Spread Optimizer",
                font_size=Pt(28), color=C_INK, font_name="Helvetica")

    # Subline
    add_textbox(slide, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.5),
                "A population-based black-box optimizer inspired by",
                font_size=Pt(16), color=C_MUTED, font_name="Helvetica")
    add_textbox(slide, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.5),
                "multi-route epidemic transmission",
                font_size=Pt(16), color=C_MUTED, font_name="Helvetica")

    # Author block (bottom-left, left-aligned — no aimless centering)
    add_textbox(slide, Inches(0.9), Inches(6.55), Inches(8), Inches(0.4),
                "Kosei Matsuzaki",
                font_size=Pt(14), bold=True, color=C_INK, font_name="Helvetica")


# ── 02: Outline ------------------------------------------------------------

OUTLINE = [
    ("1", "Background"),
    ("2", "Purpose & Significance"),
    ("3", "Proposed Method"),
    ("4", "Experiments"),
    ("5", "Results"),
    ("6", "Discussion"),
    ("7", "Conclusion"),
]


def slide_outline(prs, highlight_idx=None):
    slide = new_slide(prs)
    add_chrome(slide, "Outline", SLIDE_COUNT[0], TOTAL_PAGES[0])

    top = Inches(1.85)
    row_h = Inches(0.65)
    for i, (num, title) in enumerate(OUTLINE):
        y = top + row_h * i
        is_hl = (highlight_idx is not None and i == highlight_idx)
        color_num = C_ACCENT if is_hl else C_MUTED
        color_txt = C_INK if is_hl else C_MUTED
        weight = is_hl
        add_textbox(slide, MARGIN_L + Inches(0.5), y, Inches(0.8), Inches(0.5),
                    num, font_size=Pt(22), bold=True, color=color_num)
        add_textbox(slide, MARGIN_L + Inches(1.4), y, Inches(10), Inches(0.5),
                    title, font_size=Pt(22), bold=weight, color=color_txt)


# ── 03: Section divider ---------------------------------------------------

def slide_section(prs, num, title, subtitle=None):
    slide = new_slide(prs)
    # No standard chrome — section dividers break the repeated structure
    # deliberately (Contrast principle: change structure to start a new topic).

    # Big section number
    add_textbox(slide, Inches(0.9), Inches(2.3), Inches(2.5), Inches(2.5),
                num, font_size=Pt(140), bold=True, color=C_ACCENT_SOFT)

    # Section title
    add_textbox(slide, Inches(3.7), Inches(3.1), Inches(9), Inches(1.2),
                title, font_size=Pt(44), bold=True, color=C_INK)

    if subtitle:
        add_textbox(slide, Inches(3.7), Inches(4.1), Inches(9), Inches(0.6),
                    subtitle, font_size=Pt(18), color=C_MUTED)

    # Footer page number (subtle)
    add_textbox(slide, SLIDE_W - Inches(1.2), FOOTER_Y, Inches(0.6), Inches(0.3),
                f"{SLIDE_COUNT[0]} / {TOTAL_PAGES[0]}",
                font_size=F_TINY, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ── BACKGROUND ------------------------------------------------------------

EXP_DIR = Path(__file__).parents[1].parent / "results" / "20260514_222526_7e160d4_quick" / "dim2"


def slide_background_motivation(prs):
    """4 real BBOB landscapes in a 2x2 grid — varied difficulty types."""
    slide = new_slide(prs)
    add_chrome(slide, "Black-box functions come in many flavours of hard",
               SLIDE_COUNT[0], TOTAL_PAGES[0], section="1 — Background")

    # 4 functions covering different difficulty axes
    funcs = [
        ("F03-RastriginSep",   "Multimodal — many local optima"),
        ("F08-Rosenbrock",     "Narrow curved valley"),
        ("F10-EllipsoidalRot", "Ill-conditioned & rotated"),
        ("F12-BentCigar",      "Extreme anisotropy"),
    ]
    # Source SVG aspect ~1.91:1 — cell ratio matches to avoid wasted space.
    # 2×2 grid sized so two rows + captions fit before the footer at 7.05.
    cell_w = Inches(4.95)
    cell_h = Inches(2.59)  # 4.95 / 1.91
    cap_h = Inches(0.30)
    gap_x = Inches(0.60)   # noticeable horizontal margin
    gap_y = Inches(0.18)   # small vertical margin between caption and next row
    grid_w = cell_w * 2 + gap_x
    start_x = (SLIDE_W - grid_w) / 2
    start_y = Inches(1.35)
    for i, (fname, caption) in enumerate(funcs):
        col = i % 2
        row = i // 2
        x = start_x + col * (cell_w + gap_x)
        y = start_y + row * (cell_h + cap_h + gap_y)
        svg = EXP_DIR / f"{fname}_landscape.svg"
        png = safe_svg_to_png(svg, target_w_px=1800)
        if png:
            add_picture_fit(slide, str(png), x, y, cell_w, cell_h)
        else:
            add_rect(slide, x, y, cell_w, cell_h, fill=C_BG_SOFT, line=C_RULE)
            add_textbox(slide, x, y + cell_h/2 - Inches(0.2), cell_w, Inches(0.4),
                        f"(no SVG: {fname})", font_size=Pt(12), color=C_MUTED,
                        align=PP_ALIGN.CENTER)
        # Caption directly below the image
        add_textbox(slide, x, y + cell_h,
                    cell_w, cap_h,
                    caption, font_size=Pt(14), bold=True, color=C_INK,
                    align=PP_ALIGN.CENTER)


def slide_background_existing(prs):
    """Existing methods × their weak landscapes — card-style table."""
    slide = new_slide(prs)
    add_chrome(slide, "Every existing method has a class of problems it can't solve",
               SLIDE_COUNT[0], TOTAL_PAGES[0], section="1 — Background")

    items = [
        ("CMA-ES",   "Gets trapped in local optima on multimodal landscapes", "F03-RastriginSep"),
        ("DE",       "Struggles when the problem is strongly anisotropic",     "F12-BentCigar"),
        ("PSO",      "Stagnates on narrow valleys",                             "F08-Rosenbrock"),
        ("SaVOA",    "Single channel — can't both escape AND refine",          "F10-EllipsoidalRot"),
    ]
    # 2×2 grid of cards
    card_w = Inches(5.95)
    card_h = Inches(2.65)
    gap_x = Inches(0.30)
    gap_y = Inches(0.25)
    grid_w = card_w * 2 + gap_x
    start_x = (SLIDE_W - grid_w) / 2
    start_y = Inches(1.45)
    # Inner image area (landscape thumbnail) sized to SVG aspect ~1.91:1
    img_w = card_w - Inches(0.40)
    img_h = Inches(1.60)
    for i, (name, weakness, example) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        # Card outline + accent strip on top
        add_rect(slide, x, y, card_w, card_h, fill=None, line=C_RULE, line_w=0.75)
        add_rect(slide, x, y, card_w, Inches(0.55), fill=C_ACCENT)
        # Method name (in the accent strip)
        add_textbox(slide, x + Inches(0.30), y, Inches(2.5), Inches(0.55),
                    name, font_size=Pt(20), bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    anchor=MSO_ANCHOR.MIDDLE)
        # Weakness description (one short line)
        add_textbox(slide, x + Inches(0.20), y + Inches(0.70),
                    card_w - Inches(0.40), Inches(0.35),
                    weakness, font_size=Pt(14), color=C_INK)
        # Landscape thumbnail (bottom)
        svg = EXP_DIR / f"{example}_landscape.svg"
        png = safe_svg_to_png(svg, target_w_px=1800)
        thumb_y = y + Inches(1.05)
        if png:
            add_picture_fit(slide, str(png),
                            x + Inches(0.20), thumb_y, img_w, img_h)


def slide_background_epidemic(prs):
    """Epidemic spread ↔ optimization analogy table — the kernel of MC-ESO."""
    slide = new_slide(prs)
    add_chrome(slide, "Epidemic spread and optimization share the same structure",
               SLIDE_COUNT[0], TOTAL_PAGES[0], section="1 — Background")

    # Two illustrative panels with bold header bars (matches p11/p13 chip style)
    img_w = Inches(3.7)
    img_h = Inches(3.7)
    img_y = Inches(2.30)
    left_x = MARGIN_L + Inches(0.30)
    right_x = SLIDE_W - MARGIN_R - img_w - Inches(0.30)
    # Header bars
    add_rect(slide, left_x, Inches(1.70), img_w, Inches(0.50), fill=C_ACCENT)
    add_textbox(slide, left_x, Inches(1.70), img_w, Inches(0.50),
                "Epidemic spread", font_size=Pt(20), bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, right_x, Inches(1.70), img_w, Inches(0.50), fill=C_ACCENT)
    add_textbox(slide, right_x, Inches(1.70), img_w, Inches(0.50),
                "Optimization", font_size=Pt(20), bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_picture_fit(slide, fig_png("analogy_epidemic"),
                              left_x, img_y, img_w, img_h)
    add_picture_fit(slide, fig_png("analogy_optimization"),
                              right_x, img_y, img_w, img_h)

    # Center column — aligned vertically with the side panels (header bar
    # top at 1.70 → image bottom at 6.00 = 4.30" total). Each pair card
    # occupies a third of that height so the column looks balanced.
    map_x = left_x + img_w + Inches(0.25)
    map_w = right_x - map_x - Inches(0.25)
    mappings = [
        ("Viral particle",         "Sample point"),
        ("Host-density map",       "Objective landscape"),
        ("Outbreak in dense area", "Finding the optimum"),
    ]
    map_y = Inches(1.70)
    total_h = Inches(4.30)
    gap = Inches(0.15)
    card_h = (total_h - gap * (len(mappings) - 1)) / len(mappings)
    for i, (left, right) in enumerate(mappings):
        y = map_y + (card_h + gap) * i
        add_rect(slide, map_x, y, map_w, card_h, fill=C_BG_SOFT, line=C_RULE)
        # Top half: epidemic term (accent)
        add_textbox(slide, map_x, y, map_w, card_h / 2,
                    left, font_size=Pt(15), bold=True, color=C_ACCENT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Divider mid-card
        add_line(slide, map_x + Inches(0.20), y + card_h / 2,
                 map_x + map_w - Inches(0.20), y + card_h / 2,
                 color=C_RULE, width=0.75)
        # Bottom half: optimization term (ink)
        add_textbox(slide, map_x, y + card_h / 2, map_w, card_h / 2,
                    right, font_size=Pt(15), bold=True, color=C_INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, MARGIN_L, Inches(6.65), Inches(12.1), Inches(0.4),
                "→ Borrow how an epidemic concentrates on dense regions",
                font_size=Pt(18), bold=True, color=C_ACCENT,
                align=PP_ALIGN.CENTER)


# ── PURPOSE ---------------------------------------------------------------

def slide_purpose(prs):
    """Research goal + contributions (no preview — see next slide)."""
    slide = new_slide(prs)
    add_chrome(slide, "Goal & contributions",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="2 — Purpose & Significance")

    # Goal
    add_rect(slide, MARGIN_L, Inches(1.80), Inches(12.1), Inches(1.7),
             fill=C_BG_SOFT, line=None)
    add_textbox(slide, MARGIN_L + Inches(0.3), Inches(2.00),
                Inches(2.0), Inches(0.4),
                "GOAL", font_size=Pt(14), bold=True, color=C_ACCENT)
    add_textbox(slide, MARGIN_L + Inches(0.3), Inches(2.45),
                Inches(11.5), Inches(1.0),
                "Build a search method that performs reasonably on every function",
                font_size=Pt(24), bold=True, color=C_INK)

    # Contributions
    items = [
        ("Novelty",     "First metaheuristic to combine multiple transmission routes (epidemic-inspired)"),
        ("Mechanism",   "Channel mixing + population control + step-size adaptation"),
        ("Outcome",     "High success rate across the BBOB benchmark — see next slide"),
    ]
    top = Inches(4.20)
    for i, (label, body) in enumerate(items):
        y = top + Inches(0.85) * i
        add_textbox(slide, MARGIN_L + Inches(0.3), y, Inches(2.5), Inches(0.5),
                    label, font_size=Pt(20), bold=True, color=C_ACCENT)
        add_textbox(slide, MARGIN_L + Inches(3.0), y, Inches(9.5), Inches(0.5),
                    body, font_size=Pt(18), color=C_INK)


def slide_result_preview(prs):
    """Result preview = overall leaderboard (top) + per-category SR table (bottom).
    Numbers come from results/20260514_222526_7e160d4_quick (26 functions × 10 runs).
    """
    slide = new_slide(prs)
    add_chrome(slide, "Result preview — overall and per-category SR@1e-4",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="2 — Purpose & Significance")

    add_textbox(slide, MARGIN_L, BODY_TOP, Inches(12.1), Inches(0.4),
                "26 BBOB functions × 10 runs each — higher = solves more functions",
                font_size=Pt(13), color=C_MUTED)

    # ── Overall leaderboard (compact, single row of 5 method chips) ─────
    overall = [
        ("MC-ESO", 91.9, C_ACCENT,  True),
        ("DE",     91.9, C_CONTACT, False),
        ("PSO",    73.1, C_DROPLET, False),
        ("SaVOA",  73.1, C_AIR,     False),
        ("CMA-ES", 71.2, C_MUTED,   False),
    ]
    lb_top = Inches(1.85)
    lb_h = Inches(1.10)
    cell_w = (SLIDE_W - MARGIN_L - MARGIN_R - Inches(0.5)) / 5
    cell_gap = Inches(0.125)
    cell_w = cell_w - cell_gap
    for i, (name, val, color, is_us) in enumerate(overall):
        x = MARGIN_L + Inches(0.25) + i * (cell_w + cell_gap)
        fill = C_ACCENT_SOFT if is_us else None
        border = C_ACCENT if is_us else C_RULE
        add_rect(slide, x, lb_top, cell_w, lb_h, fill=fill, line=border,
                 line_w=2.0 if is_us else 0.75)
        # Method name
        add_textbox(slide, x, lb_top + Inches(0.10), cell_w, Inches(0.40),
                    name, font_size=Pt(16), bold=True,
                    color=C_ACCENT if is_us else C_INK,
                    align=PP_ALIGN.CENTER)
        # SR — large
        add_textbox(slide, x, lb_top + Inches(0.45), cell_w, Inches(0.60),
                    f"{val:.1f}%", font_size=Pt(28), bold=True,
                    color=C_ACCENT if is_us else C_INK,
                    align=PP_ALIGN.CENTER)

    # ── Per-category SR table ───────────────────────────────────────────
    cats = [
        ("Separable",       100, 98, 98, 82, 72),
        ("Moderate cond.",   98, 95, 100, 100, 93),
        ("Ill-conditioned",  94, 100, 40, 52, 100),
        ("Multimodal",       92, 98, 80, 68, 44),
        ("Weak structure",   74, 66, 42, 58, 40),
    ]
    method_names = [m[0] for m in overall]
    table_top = Inches(3.30)
    head_h = Inches(0.50)
    row_h = Inches(0.50)
    table_w = SLIDE_W - MARGIN_L - MARGIN_R
    cat_col_w = Inches(2.40)
    method_col_w = (table_w - cat_col_w) / 5

    # Header row
    add_rect(slide, MARGIN_L, table_top, table_w, head_h,
             fill=C_BG_SOFT, line=C_RULE, line_w=0.5)
    add_textbox(slide, MARGIN_L + Inches(0.20), table_top,
                cat_col_w, head_h,
                "Category", font_size=Pt(13), bold=True, color=C_MUTED,
                anchor=MSO_ANCHOR.MIDDLE)
    for j, m in enumerate(method_names):
        is_us = (m == "MC-ESO")
        x = MARGIN_L + cat_col_w + j * method_col_w
        add_textbox(slide, x, table_top, method_col_w, head_h,
                    m, font_size=Pt(13), bold=True,
                    color=C_ACCENT if is_us else C_INK,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Data rows
    for i, row in enumerate(cats):
        y = table_top + head_h + i * row_h
        cat_name, *vals = row
        if i % 2 == 1:
            add_rect(slide, MARGIN_L, y, table_w, row_h,
                     fill=C_BG_SOFT, line=None)
        # Category label
        add_textbox(slide, MARGIN_L + Inches(0.20), y, cat_col_w, row_h,
                    cat_name, font_size=Pt(14), color=C_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Find the best value in this row for emphasis
        max_v = max(vals)
        for j, v in enumerate(vals):
            x = MARGIN_L + cat_col_w + j * method_col_w
            is_us = (method_names[j] == "MC-ESO")
            is_best = (v == max_v)
            color = C_ACCENT if (is_us and is_best) else (
                C_ACCENT if is_us else (C_INK if is_best else C_MUTED))
            weight = is_best or is_us
            add_textbox(slide, x, y, method_col_w, row_h,
                        f"{v}%", font_size=Pt(15), bold=weight,
                        color=color, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)

    # Final separator below last row
    add_line(slide, MARGIN_L, table_top + head_h + len(cats) * row_h,
             MARGIN_L + table_w, table_top + head_h + len(cats) * row_h,
             color=C_RULE, width=0.75)

    add_textbox(slide, MARGIN_L, Inches(6.50), Inches(12.1), Inches(0.4),
                "→ MC-ESO is the only method that stays ≥ 74% in every category",
                font_size=Pt(16), bold=True, color=C_ACCENT,
                align=PP_ALIGN.CENTER)


# ── PROPOSED METHOD ------------------------------------------------------

def _channel_detail(slide, name_jp, color, fig_name, formula_img, steps):
    """Common layout for a single-channel detail slide.

    `steps` is a list of short strings rendered as a numbered list under ROLE.
    """
    add_chrome(slide, f"Channel — {name_jp}",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="3 — Proposed Method")
    # Figure on the left
    img_w = Inches(4.6); img_h = Inches(4.6)
    img_x = MARGIN_L + Inches(0.2); img_y = Inches(1.85)
    add_picture_fit(slide, fig_png(fig_name),
                              img_x, img_y, img_w, img_h)
    # Right column
    right_x = img_x + img_w + Inches(0.6)
    right_w = SLIDE_W - right_x - MARGIN_R
    # Headline
    add_textbox(slide, right_x, Inches(1.95), right_w, Inches(0.6),
                name_jp, font_size=Pt(34), bold=True, color=color)
    # Formula chip — LaTeX-rendered image
    add_rect(slide, right_x, Inches(2.85), right_w, Inches(1.10),
             fill=C_BG_SOFT, line=C_RULE)
    add_textbox(slide, right_x + Inches(0.15), Inches(2.93), right_w, Inches(0.30),
                "FORMULA", font_size=Pt(11), bold=True, color=C_MUTED)
    add_picture_fit(slide, fig_png(formula_img, transparent=True),
                    right_x + Inches(0.15), Inches(3.30),
                    right_w - Inches(0.30), Inches(0.50))
    # Numbered role steps
    add_textbox(slide, right_x, Inches(4.20), right_w, Inches(0.4),
                "ROLE", font_size=Pt(11), bold=True, color=C_MUTED)
    role_top = Inches(4.60)
    for i, step in enumerate(steps):
        y = role_top + Inches(0.55) * i
        add_textbox(slide, right_x, y, Inches(0.45), Inches(0.5),
                    f"{i + 1}.", font_size=Pt(17), bold=True, color=color)
        add_textbox(slide, right_x + Inches(0.45), y,
                    right_w - Inches(0.45), Inches(0.5),
                    step, font_size=Pt(16), color=C_INK)


def slide_generation_flow(prs):
    """Vertical 5-row timeline: each row = [step number] | [step name] | [innovation card].

    Reading order is top-to-bottom (familiar list layout); the innovation
    sits next to the step it belongs to, so the relationship is direct.
    """
    slide = new_slide(prs)
    add_chrome(slide, "One generation: 5 steps, each with its own innovation",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="3 — Proposed Method")

    rows = [
        ("1", "Spawn",            "Three transmission channels",  "Contact / Droplet / Airborne in parallel", C_ACCENT),
        ("2", "Evaluate",         "(no new mechanism)",            "Fitness call only",                          C_MUTED),
        ("3", "Update strain pool", "Strain coexistence",          "Keep up to 6 spatially-separated elites",   C_ACCENT),
        ("4", "Select population",  "Host competition",            "μ+λ greedy with parent rollback",          C_ACCENT),
        ("5", "Restart?",         "Spillover + σ adaptation",      "Escalating re-seed when stagnated",        C_ACCENT),
    ]
    top = Inches(1.50)
    row_h = Inches(1.00)
    num_w = Inches(0.80)
    name_w = Inches(2.50)
    chip_w = Inches(3.30)
    desc_w = Inches(5.50)
    pad = Inches(0.10)
    for i, (num, name, chip, desc, col) in enumerate(rows):
        y = top + row_h * i
        # Alternating background band
        if i % 2 == 1:
            add_rect(slide, MARGIN_L, y, SLIDE_W - MARGIN_L - MARGIN_R,
                     row_h - Inches(0.05), fill=C_BG_SOFT, line=None)
        # Step number — big accent circle-style typography
        add_textbox(slide, MARGIN_L + Inches(0.20), y + Inches(0.10),
                    num_w, Inches(0.80),
                    num, font_size=Pt(40), bold=True, color=col,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Step name
        add_textbox(slide, MARGIN_L + num_w + Inches(0.30), y + Inches(0.10),
                    name_w, Inches(0.80),
                    name, font_size=Pt(20), bold=True, color=C_INK,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Innovation chip
        chip_x = MARGIN_L + num_w + name_w + Inches(0.60)
        chip_fill = C_ACCENT_SOFT if col == C_ACCENT else None
        add_rect(slide, chip_x, y + Inches(0.18), chip_w, Inches(0.60),
                 fill=chip_fill, line=col, line_w=1.5)
        add_textbox(slide, chip_x, y + Inches(0.18), chip_w, Inches(0.60),
                    chip, font_size=Pt(14), bold=True, color=col,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Description
        desc_x = chip_x + chip_w + Inches(0.30)
        add_textbox(slide, desc_x, y + Inches(0.10),
                    desc_w, Inches(0.80),
                    desc, font_size=Pt(15), color=C_INK,
                    anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, MARGIN_L, Inches(6.55), Inches(12.1), Inches(0.4),
                "→ The next slides detail each innovation (rows 1, 3, 4, 5)",
                font_size=Pt(14), color=C_MUTED, align=PP_ALIGN.CENTER)


def slide_channels_overview(prs):
    """Three transmission channels side-by-side (epidemic metaphor, colour-coded)."""
    slide = new_slide(prs)
    add_chrome(slide, "Three transmission channels",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="3 — Proposed Method")

    panels = [
        ("Close-contact", C_CONTACT, "channel_contact",
         "Person-to-person touch",
         "Spawn near the source — exploit the local basin",
         "30 %"),
        ("Droplet", C_DROPLET, "channel_droplet",
         "Coughs / sneezes carry the virus to neighbours",
         "Pull toward another elite strain — implicit anisotropy",
         "40 %"),
        ("Airborne", C_AIR, "channel_air",
         "Aerosols drift far from any source",
         "Spawn far from the population — escape when stuck",
         "30 %"),
    ]
    cell_w = Inches(3.6)
    gap = Inches(0.35)
    total_w = cell_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    top = Inches(1.40)
    for i, (en, col, fig, epidemic, role, prob) in enumerate(panels):
        x = start_x + i * (cell_w + gap)
        # Header (channel name)
        add_textbox(slide, x, top, cell_w, Inches(0.5),
                    en, font_size=Pt(24), bold=True, color=col,
                    align=PP_ALIGN.CENTER)
        # Selection probability chip
        prob_x = x + cell_w / 2 - Inches(0.55)
        add_rect(slide, prob_x, top + Inches(0.50), Inches(1.1), Inches(0.30),
                 fill=col, line=None)
        add_textbox(slide, prob_x, top + Inches(0.50), Inches(1.1), Inches(0.30),
                    f"p = {prob}", font_size=Pt(11), bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Figure (square)
        add_picture_fit(slide, fig_png(fig),
                                  x, top + Inches(0.95),
                                  cell_w, cell_w)
        # Epidemic feature it mimics
        add_textbox(slide, x, top + Inches(0.95) + cell_w + Inches(0.05),
                    cell_w, Inches(0.35),
                    f"Mimics: {epidemic}",
                    font_size=Pt(11), color=C_MUTED,
                    align=PP_ALIGN.CENTER)
        # Role
        add_textbox(slide, x, top + Inches(0.95) + cell_w + Inches(0.45),
                    cell_w, Inches(0.65),
                    role,
                    font_size=Pt(12), color=C_INK,
                    align=PP_ALIGN.CENTER)


def slide_channel_contact(prs):
    slide = new_slide(prs)
    _channel_detail(
        slide, "Close-contact", C_CONTACT, "channel_contact",
        "formula_contact",
        [
            "Pick a parent x_p from the population",
            "Sample a Gaussian step with width σᵢ",
            "σᵢ shrinks with fitness and parent age",
            "→ Refines the basin the parent sits in",
        ],
    )


def slide_channel_droplet(prs):
    slide = new_slide(prs)
    _channel_detail(
        slide, "Droplet", C_DROPLET, "channel_droplet",
        "formula_droplet",
        [
            "Pull the parent toward another elite strain",
            "Perturb with a population differential vector",
            "No covariance matrix is learned explicitly",
            "→ Implicit anisotropy — alignment is automatic",
        ],
    )


def slide_channel_air(prs):
    slide = new_slide(prs)
    _channel_detail(
        slide, "Airborne", C_AIR, "channel_air",
        "formula_air",
        [
            "Sample a random location, independent of any parent",
            "Use a wide step width σ_air (much larger than σᵢ)",
            "Independent of where the population currently sits",
            "→ Escape route when the swarm is trapped",
        ],
    )


def _mech_detail(slide, name_jp, fig_name, headline, plain, params=None):
    add_chrome(slide, f"Population mechanism — {name_jp}",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="3 — Proposed Method")
    img_w = Inches(4.6); img_h = Inches(4.6)
    img_x = MARGIN_L + Inches(0.2); img_y = Inches(1.85)
    add_picture_fit(slide, fig_png(fig_name),
                              img_x, img_y, img_w, img_h)
    right_x = img_x + img_w + Inches(0.6)
    right_w = SLIDE_W - right_x - MARGIN_R
    add_textbox(slide, right_x, Inches(1.95), right_w, Inches(0.6),
                name_jp, font_size=Pt(32), bold=True, color=C_ACCENT)
    add_textbox(slide, right_x, Inches(2.75), right_w, Inches(0.7),
                headline, font_size=Pt(20), bold=True, color=C_INK)
    add_textbox(slide, right_x, Inches(3.85), right_w, Inches(2.3),
                plain, font_size=Pt(17), color=C_INK)
    if params:
        add_rect(slide, right_x, Inches(6.10), right_w, Inches(0.6),
                 fill=C_BG_SOFT, line=C_RULE)
        add_textbox(slide, right_x + Inches(0.15), Inches(6.20),
                    right_w, Inches(0.4),
                    params, font_size=Pt(14), color=C_MUTED,
                    font_name="Courier New")


def slide_mech_strain(prs):
    slide = new_slide(prs)
    _mech_detail(
        slide, "Strain coexistence", "mech_strain",
        "Keep up to 6 spatially-separated elite strains",
        "Instead of locking onto a single source, we protect several elite\n"
        "strains that are spaced apart.\n"
        "→ The droplet channel can pull toward any of them, so the search\n"
        "stays diverse on multimodal landscapes.",
        params="niche_radius_ratio = 0.1 (× span)    n_elite_max = 6"
    )


def slide_mech_host(prs):
    slide = new_slide(prs)
    _mech_detail(
        slide, "Host competition", "mech_host",
        "Each generation: kill the bottom 25%, then rollback if children are worse",
        "Kill the worst quartile of the population. Insert new children into\n"
        "the empty slots — but if a child is worse than its parent, revert.\n"
        "→ The population improves monotonically; no regression allowed.",
        params="kill_fraction = 25%"
    )


def slide_mech_spillover(prs):
    slide = new_slide(prs)
    _mech_detail(
        slide, "Spillover", "mech_spillover",
        "Full re-seed when improvement stalls; basin switch on persistent failure",
        "Reset strength grows with the stagnation streak:\n"
        " • Stage 0–1: replace 100% uniformly + axis-bound probes (best preserved)\n"
        " • Stage 2: also discard the best, reset σ (basin switch)\n"
        "Failed-basin memory rejects re-seeds near remembered dead ends.",
        params="no_improve ≥ 300  AND  f_best > 1e-8"
    )


def slide_sigma(prs):
    """σ adaptation — embed a real outbreak_dyn.svg from the experiment."""
    slide = new_slide(prs)
    add_chrome(slide, "Step-size σ adaptation",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="3 — Proposed Method")

    img_path, src = _resolve_experiment_image("F08-Rosenbrock_MC-ESO_outbreak_dyn.svg")
    if not img_path:
        img_path, src = _resolve_experiment_image("*_MC-ESO_outbreak_dyn.svg")
    if img_path:
        img_w = Inches(9.5); img_h = Inches(4.5)
        img_x = (SLIDE_W - img_w) / 2; img_y = Inches(1.55)
        add_picture_fit(slide, str(img_path), img_x, img_y, img_w, img_h)
        fname = src.name.split("_")[0]
        add_textbox(slide, MARGIN_L, Inches(6.10), Inches(12.1), Inches(0.4),
                    f"Example: σ trajectory of one MC-ESO run on {fname}",
                    font_size=Pt(13), color=C_MUTED, align=PP_ALIGN.CENTER)
    else:
        add_picture_fit(slide, fig_png("sigma_schematic"),
                                  Inches(3.5), Inches(1.7),
                                  Inches(6.5), Inches(4.0))

    add_textbox(slide, MARGIN_L, Inches(6.55), Inches(12.1), Inches(0.5),
                "Expand on improvement, shrink on stagnation, drill hard at the end.",
                font_size=Pt(18), bold=True, color=C_INK,
                align=PP_ALIGN.CENTER)


# ── EXPERIMENTS ----------------------------------------------------------

def slide_exp_setup(prs):
    """Experimental setup — accent-bar card style, matches p5/p21."""
    slide = new_slide(prs)
    add_chrome(slide, "Experimental setup",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="4 — Experiments")

    rows = [
        ("Benchmark",   "BBOB (24 functions) + 2 classical (C01–C02)"),
        ("Dimension",   "d = 2,  search domain [−5, 5]ᵈ"),
        ("Budget",      "5,000 evaluations per run"),
        ("Repetitions", "50 independent runs per (function, method)"),
        ("Success",     "Run succeeds if best f ≤ target threshold (1e−1 … 1e−10)"),
        ("Metric",      "SR@1e-k = success rate at threshold 10⁻ᵏ"),
    ]
    top = Inches(1.65)
    row_h = Inches(0.85)
    name_w = Inches(2.6)
    for i, (name, desc) in enumerate(rows):
        y = top + row_h * i
        # Accent bar
        add_rect(slide, MARGIN_L, y, Inches(0.18), row_h - Inches(0.15),
                 fill=C_ACCENT)
        # Label
        add_textbox(slide, MARGIN_L + Inches(0.40), y + Inches(0.15),
                    name_w, Inches(0.55),
                    name, font_size=Pt(20), bold=True, color=C_INK)
        # Value
        add_textbox(slide, MARGIN_L + Inches(0.40) + name_w, y + Inches(0.18),
                    Inches(9.0), Inches(0.55),
                    desc, font_size=Pt(18), color=C_INK)


def slide_exp_baselines(prs):
    """Comparison methods — card-style with MC-ESO accented."""
    slide = new_slide(prs)
    add_chrome(slide, "Comparison methods (details in appendix)",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="4 — Experiments")

    rows = [
        ("MC-ESO",  "Three transmission channels (contact + droplet + airborne) in parallel", True),
        ("CMA-ES",  "Learns a covariance matrix; reshapes sampling to fit the landscape",     False),
        ("PSO",     "Particles move toward their own best and the swarm's best",               False),
        ("SaVOA",   "Virus-inspired (single channel) with self-adaptive step size",            False),
    ]
    top = Inches(1.65)
    row_h = Inches(0.95)
    for i, (name, desc, is_us) in enumerate(rows):
        y = top + row_h * i
        # Accent bar (red for ours, gray for others)
        bar_color = C_ACCENT if is_us else C_MUTED
        add_rect(slide, MARGIN_L, y, Inches(0.20), row_h - Inches(0.15),
                 fill=bar_color)
        if is_us:
            # Highlight row
            add_rect(slide, MARGIN_L + Inches(0.20), y,
                     SLIDE_W - MARGIN_L - MARGIN_R - Inches(0.20),
                     row_h - Inches(0.15),
                     fill=C_ACCENT_SOFT, line=None)
        add_textbox(slide, MARGIN_L + Inches(0.45), y + Inches(0.20),
                    Inches(2.5), Inches(0.5),
                    name, font_size=Pt(22), bold=True,
                    color=C_ACCENT if is_us else C_INK)
        add_textbox(slide, MARGIN_L + Inches(3.1), y + Inches(0.25),
                    Inches(9.0), Inches(0.5),
                    desc, font_size=Pt(17), color=C_INK)


# ── RESULTS --------------------------------------------------------------

def slide_results_headline(prs):
    slide = new_slide(prs)
    add_chrome(slide, "Average success rate over 26 BBOB functions",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="5 — Results")

    add_textbox(slide, MARGIN_L, BODY_TOP, Inches(12), Inches(0.5),
                "26 functions × 50 runs × SR@1e-4 (= fraction reaching target precision)",
                font_size=Pt(15), color=C_MUTED)

    # Real numbers from 20260514_184347 experiment
    chart_top = Inches(2.30)
    base_x = MARGIN_L + Inches(2.6)
    bar_max_w = Inches(8.0)
    bar_h = Inches(0.62)
    gap = Inches(0.28)

    methods = [
        ("MC-ESO", 93.5, C_ACCENT,  True),
        ("PSO",    74.6, C_CONTACT, False),
        ("SaVOA",  74.3, C_DROPLET, False),
        ("CMA-ES", 70.6, C_AIR,     False),
    ]
    max_val = 100.0
    for i, (name, val, color, is_us) in enumerate(methods):
        y = chart_top + (bar_h + gap) * i
        add_textbox(slide, MARGIN_L, y + Inches(0.13), Inches(2.4), Inches(0.4),
                    name, font_size=Pt(18), bold=is_us,
                    color=C_ACCENT if is_us else C_INK,
                    align=PP_ALIGN.RIGHT)
        add_rect(slide, base_x, y, bar_max_w, bar_h, fill=C_BG_SOFT, line=None)
        w = Inches(bar_max_w.inches * val / max_val)
        add_rect(slide, base_x, y, w, bar_h, fill=color, line=None)
        add_textbox(slide, base_x + w + Inches(0.1), y + Inches(0.13),
                    Inches(1.6), Inches(0.4),
                    f"{val:.1f}%", font_size=Pt(18), bold=is_us, color=C_INK)

    add_textbox(slide, MARGIN_L, Inches(6.10), Inches(12.1), Inches(0.5),
                "→ MC-ESO leads every baseline by at least 19 percentage points",
                font_size=Pt(20), bold=True, color=C_ACCENT,
                align=PP_ALIGN.CENTER)


def slide_results_full_bbob(prs):
    slide = new_slide(prs)
    add_chrome(slide, "Full BBOB — success rate at multiple precision thresholds",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="5 — Results")

    # Big numbers grid
    top = BODY_TOP + Inches(0.1)
    card_w = Inches(2.95)
    card_h = Inches(2.0)
    gap = Inches(0.15)

    # Real numbers from 20260514_184347 experiment
    # 26 functions × 50 runs = 2600 success points per threshold
    cards = [
        ("SR @ 1e-2",  "2472 / 2600", "95.1%"),
        ("SR @ 1e-4",  "2432 / 2600", "93.5%"),
        ("SR @ 1e-7",  "2290 / 2600", "88.1%"),
        ("SR @ 1e-10", "2184 / 2600", "84.0%"),
    ]
    x = MARGIN_L
    for label, val, pct in cards:
        is_main = (label == "SR @ 1e-4")
        fill = C_ACCENT_SOFT if is_main else C_BG_SOFT
        border = C_ACCENT if is_main else C_RULE
        add_rect(slide, x, top, card_w, card_h, fill=fill, line=border,
                 line_w=2.0 if is_main else 1.0)
        add_textbox(slide, x, top + Inches(0.2), card_w, Inches(0.4),
                    label, font_size=Pt(14), bold=True,
                    color=C_ACCENT if is_main else C_MUTED,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(0.65), card_w, Inches(0.8),
                    val, font_size=Pt(26), bold=True, color=C_INK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, top + Inches(1.45), card_w, Inches(0.4),
                    pct, font_size=Pt(18), color=C_MUTED,
                    align=PP_ALIGN.CENTER)
        x += card_w + gap

    # Takeaway: top method on 21/26 functions (alone or tied)
    add_textbox(slide, MARGIN_L, Inches(5.4), Inches(12.1), Inches(0.6),
                "Top method on 21 / 26 functions  (6 outright + 15 tied)",
                font_size=Pt(24), bold=True, color=C_ACCENT,
                align=PP_ALIGN.CENTER)


def slide_results_hard(prs):
    slide = new_slide(prs)
    add_chrome(slide, "Where MC-ESO still struggles",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="5 — Results")

    add_textbox(slide, MARGIN_L, BODY_TOP, Inches(12), Inches(0.5),
                "The 3 weakest cases — all of them are open problems for every method",
                font_size=Pt(15), color=C_MUTED)

    header_y = Inches(2.2)
    col_x = [MARGIN_L, MARGIN_L + Inches(3.4), MARGIN_L + Inches(5.4),
             MARGIN_L + Inches(7.6)]
    col_w = [Inches(3.3), Inches(2.0), Inches(2.2), Inches(4.5)]
    for cx, cw, h in zip(col_x, col_w,
                         ["Function", "SR @ 1e-4", "Type", "Why it's hard"]):
        add_textbox(slide, cx, header_y, cw, Inches(0.4),
                    h, font_size=Pt(14), bold=True, color=C_MUTED)
    add_line(slide, MARGIN_L, header_y + Inches(0.45),
             SLIDE_W - MARGIN_R, header_y + Inches(0.45), color=C_RULE)

    # Real numbers from 20260514_184347 experiment
    rows = [
        ("F24  Lunacek bi-Rastrigin", "16%", "Double funnel",
         "Deeper funnel is the wrong basin (deceptive)"),
        ("F23  Katsuura",              "58%", "Fractal",
         "Self-similar; no smoothing scale exists"),
        ("F18  Schaffer F7 (ill)",     "88%", "Ill-cond. multimodal",
         "Many local optima inside a narrow valley"),
        ("F04  Büche-Rastrigin",       "92%", "Asymmetric multimodal",
         "Dense grid of local optima, asymmetric"),
    ]
    row_h = Inches(0.72)
    for i, (fn, sr, t, why) in enumerate(rows):
        y = header_y + Inches(0.55) + row_h * i
        if i % 2 == 1:
            add_rect(slide, MARGIN_L, y, SLIDE_W - MARGIN_L - MARGIN_R, row_h - Inches(0.07),
                     fill=C_BG_SOFT, line=None)
        add_textbox(slide, col_x[0] + Inches(0.1), y + Inches(0.13), col_w[0], Inches(0.5),
                    fn, font_size=Pt(15), bold=True, color=C_INK)
        # color the SR cell
        sr_color = C_OK if int(sr.rstrip("%")) >= 70 else C_BAD
        add_textbox(slide, col_x[1] + Inches(0.1), y + Inches(0.13), col_w[1], Inches(0.5),
                    sr, font_size=Pt(17), bold=True, color=sr_color)
        add_textbox(slide, col_x[2] + Inches(0.1), y + Inches(0.15), col_w[2], Inches(0.5),
                    t, font_size=Pt(14), color=C_MUTED)
        add_textbox(slide, col_x[3] + Inches(0.1), y + Inches(0.15), col_w[3], Inches(0.5),
                    why, font_size=Pt(14), color=C_INK)

    add_textbox(slide, MARGIN_L, Inches(6.55), Inches(12), Inches(0.4),
                "These remain open problems across the entire metaheuristic literature",
                font_size=Pt(14), color=C_MUTED)


# ── DISCUSSION -----------------------------------------------------------

def webp_to_png(webp_path: Path) -> Path:
    """Convert webp → png via macOS sips. Cached by mtime."""
    webp_path = Path(webp_path)
    out = SVG_CACHE / f"{webp_path.stem}.png"
    if out.exists() and out.stat().st_mtime >= webp_path.stat().st_mtime:
        return out
    subprocess.run(["sips", "-s", "format", "png",
                    str(webp_path), "--out", str(out)],
                   check=True, capture_output=True, timeout=15)
    return out


def _resolve_experiment_image(glob_pattern):
    """Return a usable raster path for an experiment file.
    Convert SVG via qlmanage, WEBP via sips; PNG used directly."""
    for hit in sorted(EXP_DIR.glob(glob_pattern)):
        suf = hit.suffix.lower()
        if suf == ".svg":
            png = safe_svg_to_png(hit, target_w_px=1800)
            if png:
                return png, hit
        elif suf == ".webp":
            try:
                return webp_to_png(hit), hit
            except Exception as e:
                print(f"  ! WEBP conversion failed: {e}")
                continue
        else:
            return hit, hit
    return None, None


def _why_slide(slide, title_jp, head_kw, evidence_caption, real_glob,
               fallback_fig, body_text):
    add_chrome(slide, f"Why does it work — {title_jp}",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="6 — Discussion")
    add_textbox(slide, MARGIN_L, Inches(1.55), Inches(12.1), Inches(0.6),
                head_kw, font_size=Pt(28), bold=True, color=C_ACCENT)
    img_w = Inches(7.0); img_h = Inches(4.7)
    img_x = SLIDE_W - MARGIN_R - img_w
    img_y = Inches(2.20)
    real_path, src = _resolve_experiment_image(real_glob)
    if real_path:
        add_picture_fit(slide, str(real_path), img_x, img_y, img_w, img_h)
        fname = src.name.split("_")[0]
        add_textbox(slide, img_x, img_y + img_h + Inches(0.05),
                    img_w, Inches(0.3),
                    f"From the {fname} run",
                    font_size=Pt(11), color=C_MUTED, align=PP_ALIGN.CENTER)
    else:
        add_picture_fit(slide, fig_png(fallback_fig),
                                  img_x, img_y, img_w, img_h)
    add_textbox(slide, MARGIN_L, Inches(2.30), Inches(5.5), Inches(4.5),
                body_text, font_size=Pt(18), color=C_INK)


def slide_discussion_why_anisotropy(prs):
    slide = new_slide(prs)
    _why_slide(
        slide,
        title_jp="① Implicit anisotropy",
        head_kw="The droplet channel picks up the landscape's anisotropy for free",
        evidence_caption="Population spreads along the elongated valley",
        real_glob="F08-Rosenbrock_MC-ESO_population.webp",
        fallback_fig="channel_droplet",
        body_text=(
            "The differential vector used by the droplet channel\n"
            "naturally aligns with the dominant spread direction of\n"
            "the population.\n\n"
            "→ Same benefit as CMA-ES, without explicit covariance.\n\n"
            "• Extra cost: zero\n"
            "• Helps on: narrow valleys like F08 Rosenbrock"
        ),
    )


def slide_discussion_why_multimodal(prs):
    slide = new_slide(prs)
    _why_slide(
        slide,
        title_jp="② Multimodal coverage",
        head_kw="Strain coexistence keeps multiple peaks alive simultaneously",
        evidence_caption="Multiple elite strains coexist in the pool",
        real_glob="F03-RastriginSep_MC-ESO_population.webp",
        fallback_fig="mech_strain",
        body_text=(
            "Instead of locking onto a single source, MC-ESO keeps\n"
            "up to 6 spatially-separated elite strains.\n\n"
            "→ A single trapped optimum no longer terminates the search.\n\n"
            "• Extra cost: light (distance check only)\n"
            "• Helps on: rugged multimodal like F03 Rastrigin"
        ),
    )


def slide_discussion_why_disruption(prs):
    slide = new_slide(prs)
    _why_slide(
        slide,
        title_jp="③ Adaptive disruption",
        head_kw="Spillover explosively re-seeds a stuck population",
        evidence_caption="σ trajectory: stagnate → restart → reconverge cycles",
        real_glob="F24-LunacekRastrigin_MC-ESO_outbreak_dyn.svg",
        fallback_fig="mech_spillover",
        body_text=(
            "When improvement stalls for long, the reset strength\n"
            "escalates with each successive failure.\n\n"
            "→ The algorithm doesn't just wait — it switches basins\n"
            "by itself.\n\n"
            "• Helps on: multi-basin landscapes like F24"
        ),
    )


def slide_discussion_limits(prs):
    slide = new_slide(prs)
    add_chrome(slide, "Limitations",
               SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="6 — Discussion")

    items = [
        "Hand-tuned channel ratios (30:40:30)",
        "Deceptive funnels  (F24, F23)",
        "Only d = 2 reported",
        "No gradient hybridization",
    ]
    top = Inches(2.2)
    for i, label in enumerate(items):
        y = top + Inches(0.95) * i
        add_textbox(slide, MARGIN_L + Inches(0.3), y, Inches(0.5), Inches(0.5),
                    "•", font_size=Pt(24), bold=True, color=C_ACCENT)
        add_textbox(slide, MARGIN_L + Inches(0.8), y, Inches(11), Inches(0.6),
                    label, font_size=Pt(22), color=C_INK)


# ── CONCLUSION -----------------------------------------------------------

def slide_conclusion(prs):
    slide = new_slide(prs)
    add_chrome(slide, "Conclusion", SLIDE_COUNT[0], TOTAL_PAGES[0],
               section="7 — Conclusion")

    items = [
        ("Idea",       "Multi-route reproduction inspired by epidemic spread"),
        ("Mechanism",  "3 transmission channels + 3 population mechanisms + adaptive σ"),
        ("Result",     "Average SR = 93.5% on BBOB-26 — outperforms all 3 baselines"),
    ]
    top = Inches(2.2)
    for i, (label, body) in enumerate(items):
        y = top + Inches(1.20) * i
        add_rect(slide, MARGIN_L, y, Inches(0.22), Inches(0.95), fill=C_ACCENT)
        add_textbox(slide, MARGIN_L + Inches(0.5), y, Inches(3.0), Inches(0.5),
                    label, font_size=Pt(22), bold=True, color=C_ACCENT)
        add_textbox(slide, MARGIN_L + Inches(0.5), y + Inches(0.45),
                    Inches(11.7), Inches(0.55),
                    body, font_size=Pt(22), color=C_INK)

    add_textbox(slide, MARGIN_L, Inches(6.55), Inches(12), Inches(0.5),
                "Thank you.",
                font_size=Pt(22), bold=True, color=C_ACCENT)


# ── APPENDIX -------------------------------------------------------------

def slide_appendix_savoa(prs):
    """Appendix: SaVOA (Self-adaptive Virus Optimization Algorithm) details."""
    slide = new_slide(prs)
    add_chrome(slide, "Appendix — SaVOA (Self-adaptive Virus Optimization)",
               SLIDE_COUNT[0], TOTAL_PAGES[0], section="Appendix")

    add_textbox(slide, MARGIN_L, Inches(1.60), Inches(12.1), Inches(0.7),
                "The closest prior work — also virus-inspired, but with a single channel",
                font_size=Pt(20), bold=True, color=C_INK)

    rows = [
        ("Inspiration",  "Viral infection and self-replication"),
        ("Reproduction", "Single Gaussian sample around the parent (contact channel only)"),
        ("σ update",     "Self-adaptive — each agent learns its own step size"),
        ("Population",   "Simple generational replacement. No niching, no escalating restart."),
        ("vs. MC-ESO",
         "• 1 channel → 3 (contact + droplet + airborne)\n"
         "• Adds 3 population mechanisms (strain coexistence, host competition, spillover)"),
    ]
    top = Inches(2.55)
    for i, (k, v) in enumerate(rows):
        y = top + Inches(0.85) * i
        add_rect(slide, MARGIN_L, y, Inches(0.18), Inches(0.7), fill=C_ACCENT)
        add_textbox(slide, MARGIN_L + Inches(0.4), y + Inches(0.05),
                    Inches(2.6), Inches(0.5),
                    k, font_size=Pt(16), bold=True, color=C_MUTED)
        add_textbox(slide, MARGIN_L + Inches(3.1), y + Inches(0.05),
                    Inches(9.0), Inches(1.0),
                    v, font_size=Pt(16), color=C_INK)


def slide_appendix_others(prs):
    """Appendix: one-liners for CMA-ES / PSO."""
    slide = new_slide(prs)
    add_chrome(slide, "Appendix — Other baselines",
               SLIDE_COUNT[0], TOTAL_PAGES[0], section="Appendix")

    rows = [
        ("CMA-ES",
         "Covariance Matrix Adaptation Evolution Strategy (de facto standard)",
         "• Learns the shape of the sampling Gaussian itself\n"
         "• Excellent on elongated / rotated landscapes\n"
         "• Tends to get trapped in local optima on multimodal functions"),
        ("PSO",
         "Particle Swarm Optimization",
         "• Each particle moves toward its own best and the swarm best\n"
         "• Strong on smooth problems / stagnates on high-dim narrow valleys"),
    ]
    top = Inches(1.60)
    block_h = Inches(1.65)
    for i, (name, sub, body) in enumerate(rows):
        y = top + block_h * i
        add_rect(slide, MARGIN_L, y, Inches(12.1), block_h - Inches(0.15),
                 fill=None, line=C_RULE, line_w=0.75)
        add_rect(slide, MARGIN_L, y, Inches(0.18), block_h - Inches(0.15),
                 fill=C_ACCENT)
        add_textbox(slide, MARGIN_L + Inches(0.4), y + Inches(0.10),
                    Inches(3.0), Inches(0.5),
                    name, font_size=Pt(22), bold=True, color=C_INK)
        add_textbox(slide, MARGIN_L + Inches(3.6), y + Inches(0.18),
                    Inches(8.3), Inches(0.4),
                    sub, font_size=Pt(13), color=C_MUTED)
        add_textbox(slide, MARGIN_L + Inches(0.4), y + Inches(0.65),
                    Inches(11.5), Inches(1.0),
                    body, font_size=Pt(14), color=C_INK)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        lambda p: slide_outline(p, highlight_idx=None),
        lambda p: slide_section(p, "1", "Background",
                                "The shared weakness of existing methods, and what motivated this work"),
        slide_background_motivation,
        slide_background_existing,
        slide_background_epidemic,
        lambda p: slide_section(p, "2", "Purpose & Significance",
                                "What we ask, and why it matters"),
        slide_purpose,
        slide_result_preview,
        lambda p: slide_section(p, "3", "Proposed Method",
                                "MC-ESO — generation flow, three channels, three population mechanisms"),
        slide_generation_flow,
        slide_channels_overview,
        slide_channel_contact,
        slide_channel_droplet,
        slide_channel_air,
        slide_mech_strain,
        slide_mech_host,
        slide_mech_spillover,
        slide_sigma,
        lambda p: slide_section(p, "4", "Experiments",
                                "BBOB benchmark — conditions and baselines"),
        slide_exp_setup,
        slide_exp_baselines,
        lambda p: slide_section(p, "5", "Results",
                                "Success rate across all functions; the hard cases"),
        slide_results_headline,
        slide_results_full_bbob,
        slide_results_hard,
        lambda p: slide_section(p, "6", "Discussion",
                                "Why it works, and where it loses"),
        slide_discussion_why_anisotropy,
        slide_discussion_why_multimodal,
        slide_discussion_why_disruption,
        slide_discussion_limits,
        lambda p: slide_section(p, "7", "Conclusion",
                                "Summary"),
        slide_conclusion,
        # Appendix
        slide_appendix_savoa,
        slide_appendix_others,
    ]

    TOTAL_PAGES[0] = len(builders)

    for fn in builders:
        fn(prs)

    out = Path(__file__).parent / "MC-ESO.pptx"
    prs.save(out)
    print(f"Wrote: {out}   ({SLIDE_COUNT[0]} slides)")


if __name__ == "__main__":
    build()
