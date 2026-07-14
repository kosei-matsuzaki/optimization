"""Build the 2026-07-14 MC-ESO progress-report deck (#2) with python-pptx.

Reproduces the visual language of presentation/20260519.pdf: red sidebar/accent,
large section-divider numbers, one-message content slides with a title rule and
footer, stat tiles, tables, and W/T/L chips. English text to match the prior deck.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
FIG = HERE / "figs"
OUTDIR = HERE.parent
FONT = "Arial"

# ── palette ──────────────────────────────────────────────────────────────
RED     = RGBColor(0xC0, 0x39, 0x2B)
RED_DK  = RGBColor(0x7C, 0x1C, 0x11)
DARK    = RGBColor(0x1F, 0x27, 0x33)
GRAY    = RGBColor(0x5B, 0x66, 0x73)
PINK    = RGBColor(0xE8, 0xA9, 0xA0)
BG      = RGBColor(0xF2, 0xF4, 0xF6)
BORDER  = RGBColor(0xD9, 0xDC, 0xE0)
THEAD   = RGBColor(0xED, 0xEF, 0xF2)
TALT    = RGBColor(0xF7, 0xF8, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
RULEC   = RGBColor(0xD0, 0xD4, 0xD9)
BLUE    = RGBColor(0x2E, 0x6D, 0xA4)
GREEN   = RGBColor(0x4A, 0x8B, 0x5C)
TEAL    = RGBColor(0x2E, 0x8B, 0x8B)
AMBER   = RGBColor(0xCF, 0x8A, 0x2B)
PURPLE  = RGBColor(0x6B, 0x4E, 0x9E)
GRN_DK  = RGBColor(0x2F, 0x6B, 0x3E)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

TOTAL_PAGES = None  # filled after build
_PAGE_NO = 0        # actual slide ordinal, auto-incremented in slide()


def slide():
    global _PAGE_NO
    _PAGE_NO += 1
    return prs.slides.add_slide(BLANK)


def _set_run(r, size, bold, color, italic=False, name=FONT):
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    r.font.color.rgb = color


def txt(s, l, t, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, name=FONT, spacing=None):
    """Single-run textbox."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run(); r.text = text
    _set_run(r, size, bold, color, italic, name)
    return tb


def paras(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """Multi-paragraph textbox. lines = list of dict(runs=[(txt,size,bold,color,ital)],
    align, space_before, bullet, spacing)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_before") is not None:
            p.space_before = Pt(ln["space_before"])
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        if ln.get("spacing"):
            p.line_spacing = ln["spacing"]
        for (t_, sz, bd, col, *rest) in ln["runs"]:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t_
            _set_run(r, sz, bd, col, it)
    return tb


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, radius=None):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    if radius is not None:                  # corner radius of a rounded rect
        try:
            sp.adjustments[0] = radius      # fraction of the shorter side
        except (IndexError, ValueError):
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if not shadow:
        pass
    return sp


BOXR_IN = 0.08    # ABSOLUTE corner radius (inches) — identical on every box,
                  # whatever its size, so corners read the same deck-wide


def _box_frac(w, h):
    """Rounded-rect adjustment fraction that yields the constant BOXR_IN radius
    for a box of this size (python-pptx expresses the radius as a fraction of the
    shorter side, so a fixed absolute radius needs a per-box fraction)."""
    short = min(int(w), int(h))
    return min(0.5, int(Inches(BOXR_IN)) / short) if short else 0.0


BOX_LW = 1.2      # single border weight for every container box on the deck

# heading tick: the short vertical rule that sits left of a section heading.
HTICK_W = 0.07    # thickness (inches)
HTICK_H = 0.30    # height (inches)
HTICK_GAP = 0.18  # gap from the tick's right edge to the heading text (inches)


def softbox(s, l, t, w, h, fill=None, line=None, line_w=None):
    """A container box with the deck-wide soft corner radius AND border weight.
    Every callout / info panel / equation box / card routes through this so the
    corner + border treatment is identical on every page. (line_w is accepted
    for call-site compatibility but ignored — the deck uses one weight, BOX_LW.)"""
    return rect(s, l, t, w, h, fill=fill, line=line, line_w=BOX_LW,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=_box_frac(w, h))


def htick(s, x, top, color=RED, h=Inches(HTICK_H)):
    """Draw a heading tick (fixed thickness/height/gap) with its top-left at
    (x, top) — both EMU/Length — and return the x where the heading text starts."""
    rect(s, x, top, Inches(HTICK_W), h, fill=color)
    return Emu(int(x) + int(Inches(HTICK_W)) + int(Inches(HTICK_GAP)))


SPINE = RGBColor(0x9A, 0xA2, 0xAA)   # neutral structural line for branch trees


def branch_fork(s, cx, top, branches, eq_top, bus_gap=Inches(0.20)):
    """A condition box forks into labelled branches. The whole tree structure
    (stub down / one bus across / a short stub into each branch) is a single
    neutral colour; only the leaf carries the branch colour (filled pill + arrow
    into that branch's equation box). Shared by the restart (①) and best2 (④)
    mechanism slides so the two conditionals read identically.

    cx/top: centre-x and bottom-y of the condition box (Length).
    branches: [(bx, pill_w, label, colour), ...] — bx/pill_w are Length.
    eq_top: top-y of the equation boxes the arrows point into (Length)."""
    bus = Emu(int(top) + int(bus_gap))
    xs = [int(b[0]) for b in branches]
    connector(s, cx, top, cx, bus, color=SPINE, weight=1.6, arrow=False)
    connector(s, Emu(min(xs)), bus, Emu(max(xs)), bus, color=SPINE, weight=1.6,
              arrow=False)
    drop, ph = Inches(0.16), Inches(0.34)
    p_top = Emu(int(bus) + int(drop))
    p_bot = Emu(int(p_top) + int(ph))
    a_end = Emu(int(eq_top) - int(Inches(0.04)))
    for bx, pw, label, col in branches:
        connector(s, bx, bus, bx, p_top, color=SPINE, weight=1.6, arrow=False)
        px = Emu(int(bx) - int(pw) // 2)
        softbox(s, px, p_top, pw, ph, fill=col)
        txt(s, px, p_top, pw, ph, label, size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        connector(s, bx, p_bot, bx, a_end, color=col, weight=1.8)


def condition_box(s, left, top, w, h):
    """The gate at the top of a branch_fork, drawn as a stadium (fully rounded
    semicircular ends) so it reads as a condition — deliberately distinct from
    the deck's rounded-rectangle content boxes."""
    return rect(s, left, top, w, h, fill=BG, line=RED, line_w=1.4,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)


def hline(s, l, t, w, color=RULEC, weight=1.4):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, l, t, l + w, t)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def connector(s, x1, y1, x2, y2, color=RULEC, weight=1.4, arrow=True):
    """Straight connector from (x1,y1) to (x2,y2), optional triangle arrowhead."""
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    if arrow:
        el = ln.line._get_or_add_ln()
        el.append(el.makeelement(qn("a:tailEnd"), {"type": "triangle"}))
    return ln


def chrome(s, title, section, page=None):
    """Content-slide chrome: title, rule, footer, page number.

    The page number is the actual slide ordinal (auto-tracked in slide()), so
    inserting or reordering slides never requires renumbering by hand.
    """
    txt(s, MARGIN, Inches(0.32), Inches(12.2), Inches(0.72), title,
        size=30, bold=True, color=DARK)
    hline(s, MARGIN, Inches(1.04), Inches(13.333) - 2 * MARGIN)
    txt(s, MARGIN, Inches(7.02), Inches(6), Inches(0.3),
        f"{section}", size=10.5, color=GRAY)
    txt(s, Inches(11.2), Inches(7.02), Inches(1.6), Inches(0.3),
        f"{_PAGE_NO} / {{TP}}", size=10.5, color=GRAY, align=PP_ALIGN.RIGHT)


def subtitle(s, text, y=Inches(1.32)):
    txt(s, MARGIN, y, Inches(12.0), Inches(0.45), text, size=17, color=DARK)


def datanote(s, text, top):
    """Experiment-condition / reading-key note for a chart or table: always
    right-aligned directly above the element (0.36" above its top edge) —
    one fixed position deck-wide instead of ad-hoc footer captions."""
    txt(s, MARGIN, top - Inches(0.36), Inches(12.23), Inches(0.3), text,
        size=12, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════
def title_slide():
    s = slide()
    rect(s, 0, 0, Inches(0.28), SH, fill=RED)
    txt(s, Inches(0.9), Inches(2.5), Inches(8), Inches(0.4),
        "PROGRESS REPORT #2", size=15, bold=True, color=RED)
    txt(s, Inches(0.86), Inches(2.95), Inches(10), Inches(1.5),
        "MC-ESO", size=76, bold=True, color=DARK)
    txt(s, Inches(0.9), Inches(4.35), Inches(11), Inches(0.6),
        "Multi-Channel Epidemic Spread Optimizer", size=28, color=DARK)
    paras(s, Inches(0.9), Inches(5.05), Inches(11), Inches(0.9), [
        {"runs": [("Diagnosis-driven refinements since 5/18",
                   14, False, GRAY)]},
    ])
    txt(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
        "M1 Kosei Matsuzaki   /   2026-07-14", size=14, bold=True, color=DARK)


def divider(num, title, sub):
    s = slide()
    txt(s, Inches(0.5), Inches(1.7), Inches(3.2), Inches(4.2), str(num),
        size=250, bold=True, color=PINK, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.9), Inches(3.05), Inches(8.8), Inches(1.0), title,
        size=46, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.95), Inches(4.05), Inches(8.6), Inches(0.6), sub,
        size=17, color=GRAY)


def card(s, l, t, w, h, header, body_lines, header_fill=RED, body_fill=WHITE):
    """A red-header card with a white body (as in the prior deck)."""
    hh = Inches(0.62)
    rect(s, l, t, w, h, fill=body_fill, line=BORDER, line_w=1.0)
    rect(s, l, t, w, hh, fill=header_fill)
    txt(s, l + Inches(0.22), t, w - Inches(0.4), hh, header,
        size=17, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    paras(s, l + Inches(0.22), t + hh + Inches(0.14), w - Inches(0.44),
          h - hh - Inches(0.28), body_lines)


def accent_item(s, l, t, w, label, desc, label_color=RED, label_size=15, desc_size=13):
    """Red vertical tick + bold label + gray desc (prior deck's list style)."""
    rect(s, l, t + Inches(0.02), Inches(0.06), Inches(0.5), fill=label_color)
    paras(s, l + Inches(0.2), t, w - Inches(0.2), Inches(0.9), [
        {"runs": [(label, label_size, True, DARK)]},
        {"runs": [(desc, desc_size, False, GRAY)], "space_before": 2},
    ])


def chip(s, l, t, w, top, big):
    rect(s, l, t, w, Inches(0.85), fill=BG, line=BORDER, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE)
    paras(s, l, t + Inches(0.12), w, Inches(0.65), [
        {"runs": [(top, 11, False, GRAY)], "align": PP_ALIGN.CENTER},
        {"runs": [(big, 14, True, RED)], "align": PP_ALIGN.CENTER, "space_before": 3},
    ])


def table(s, l, t, col_w, rows, header_bold=True, row_h=Inches(0.44),
          first_bold=True, highlight_row=None, fontsize=13,
          cell_fills=None, cell_text=None,
          inset_l=Inches(0.12), inset_r=Inches(0.08)):
    """Lightweight table via rectangles + text. rows[0] = header.
    cell_fills/cell_text: optional {(ri, ci): color} overrides for single cells."""
    cell_fills = cell_fills or {}
    cell_text = cell_text or {}
    y = t
    for ri, row in enumerate(rows):
        x = l
        if ri == 0:
            fill = THEAD
        elif highlight_row is not None and ri == highlight_row:
            fill = PINK
        elif ri % 2 == 0:
            fill = TALT
        else:
            fill = WHITE
        total_w = sum(col_w)
        rect(s, x, y, total_w, row_h, fill=fill,
             line=BORDER if ri == 0 else None, line_w=0.75)
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            if (ri, ci) in cell_fills:
                rect(s, x, y, cw, row_h, fill=cell_fills[(ri, ci)])
            bold = (ri == 0 and header_bold) or (ci == 0 and first_bold)
            col = DARK
            if ri == 0:
                col = GRAY
            if highlight_row is not None and ri == highlight_row:
                col = RED_DK; bold = True
            if (ri, ci) in cell_text:
                col = cell_text[(ri, ci)]; bold = True
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            txt(s, x + inset_l, y, cw - inset_l - inset_r, row_h, str(cell),
                size=fontsize, bold=bold, color=col, align=align,
                anchor=MSO_ANCHOR.MIDDLE)
            x += cw
        y = Emu(int(y) + int(row_h))
    # bottom rule
    hline(s, l, y, sum(col_w), color=BORDER, weight=1.0)


def img(s, path, l, t, w=None, h=None, caption=None):
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    pic = s.shapes.add_picture(str(path), l, t, **kw)
    if caption:
        txt(s, l, Emu(int(t) + int(pic.height)) , pic.width, Inches(0.3), caption,
            size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
    return pic


# ── content slides ─────────────────────────────────────────────────────────
def p2_motivation():
    s = slide()
    chrome(s, "Motivation & goal", "1 — Recap", 2)
    subtitle(s, "Every method has a landscape it fails on")
    # method → weak-shape table (evidence from the 10-method comparison)
    col_w = [Inches(3.55), Inches(5.35), Inches(3.33)]
    rows = [
        ["Method", "Weak function shape", "SR@1e-10"],
        ["CMA-ES / L-SHADE", "Multimodal  (F15–F24)", "32% / 42%"],
        ["PSO / SaVOA", "Ill-conditioned  (F10–F14)", "7% / 18%"],
        ["DE", "Bent-valley / asymmetric  (F06–F09)", "91%  (rest ≈100)"],
    ]
    ty = Inches(2.15)
    datanote(s, "mean SR@1e-10 / BBOB-24 / dim 2", ty)
    table(s, MARGIN, ty, col_w, rows, row_h=Inches(0.66), fontsize=14.5)
    # goal band
    softbox(s, MARGIN, Inches(5.35), Inches(12.23), Inches(1.35), fill=BG,
            line=RED, line_w=1.5)
    paras(s, Inches(0.85), Inches(5.62), Inches(11.6), Inches(0.9), [
        {"runs": [("GOAL", 12, True, RED)]},
        {"runs": [("Mimic multi-route viral spread — robust across every function shape",
                   22, True, DARK)], "space_before": 6},
    ])


def _dim_panel(s, l, t, w, h, alpha=0.25):
    """Overlay a semi-transparent black rectangle to grey a panel out (used to
    highlight one schematic by dimming the others)."""
    sp = rect(s, l, t, w, h, fill=RGBColor(0x0A, 0x0E, 0x14))
    srgb = sp._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    srgb.append(srgb.makeelement(qn("a:alpha"),
                                 {"val": str(int(alpha * 100000))}))
    return sp


def p3_method_grid(highlight=None, title="How MC-ESO works",
                   section="1 — Recap"):
    """The whole method on one page as a 3x2 gallery: the top row is the three
    transmission channels (how a host reproduces), the bottom row the three
    population strategies (how the swarm is steered). Each row carries a
    labelled header on the left; the schematics sit plainly on the slide, each
    with its name in the family accent colour.

    When `highlight` is a panel name (or a set of names) the OTHER panels are
    dimmed with a semi-transparent overlay — reused before each improvement to
    show which mechanism(s) that change refines."""
    hl = ({highlight} if isinstance(highlight, str)
          else set(highlight) if highlight else None)
    s = slide()
    chrome(s, title, section)
    # Group labels move to a left gutter (instead of a full-width header band),
    # so the two schematic rows get real vertical breathing room between them.
    gx0, gx1 = 2.25, 12.78                       # panels occupy the space right of the gutter
    centers = [Inches(gx0 + (gx1 - gx0) * (i + 0.5) / 3) for i in range(3)]
    tw, iw = Inches(3.3), Inches(3.0)
    # single horizontal divider separating the two group bands
    hline(s, MARGIN, Inches(4.12), Inches(13.333) - 2 * MARGIN,
          color=RULEC, weight=1.0)
    rows = [
        # (name-row y, label-centre y, group name, gloss, panels)
        (1.62, 2.77, "Transmission\nchannels", "three parallel\nreproduction routes", [
            ("Droplet", RED, "p04_channels/droplet.emf"),
            ("Close-contact", BLUE, "p04_channels/contact.emf"),
            ("Airborne", GREEN, "p04_channels/airborne.emf")]),
        (4.48, 5.52, "Population\nstrategies", "how the swarm\nis steered", [
            ("Strain coexistence", TEAL, "p05_mechanisms/strain.emf"),
            ("Restart  (spillover)", PURPLE, "p05_mechanisms/restart.emf"),
            ("Drilling", AMBER, "p05_mechanisms/drilling.emf")]),
    ]
    for ny, lcy, head, sub, panels in rows:
        # left gutter: clean group name + gloss, vertically centred on the row's
        # schematics. One paragraph per line so spacing is fully controlled:
        # tight within the name / within the gloss, a wide gap between them.
        lines = []
        for i, hln in enumerate(head.split("\n")):
            lines.append({"runs": [(hln, 16, True, DARK)],
                          "space_before": 0 if i == 0 else 2})
        for j, sl in enumerate(sub.split("\n")):
            lines.append({"runs": [(sl, 10.5, False, GRAY)],
                          "space_before": 12 if j == 0 else 2})
        paras(s, MARGIN, Inches(lcy - 0.60), Inches(1.75), Inches(1.2), lines,
              anchor=MSO_ANCHOR.MIDDLE)
        for (nm, col, fn), cx in zip(panels, centers):
            txt(s, cx - tw / 2, Inches(ny), tw, Inches(0.34), nm,
                size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
            img(s, FIG / fn, cx - iw / 2, Inches(ny + 0.42), w=iw)
            if hl is not None and nm not in hl:
                _dim_panel(s, Emu(int(cx) - int(Inches(1.675))),
                           Inches(ny - 0.20), Inches(3.35), Inches(2.55))


def p4_prev_result():
    s = slide()
    chrome(s, "Previous results  (5/18)", "1 — Recap")
    subtitle(s, "Best at SR@1e-4 — but DE led at 1e-10")
    datanote(s, "BBOB-24 [10] / dim 2 / n = 30", Inches(2.0))
    # left: paired-bar chart (shared style with the "now" slide)
    img(s, FIG / "p06_prev_result/prev518.emf", MARGIN, Inches(2.0), w=Inches(7.95))
    # right: exact SR table; colour tells the story instead of a caption —
    # green = best in that column, red = MC-ESO behind (the deep-precision gap)
    PALE_GRN = RGBColor(0xDF, 0xEE, 0xE3)
    PALE_RED = RGBColor(0xF7, 0xE0, 0xDD)
    GRN_TX = RGBColor(0x2F, 0x6B, 0x3E)
    x = Inches(8.7)
    table(s, x, Inches(2.8), [Inches(1.8), Inches(1.15), Inches(1.15)], [
        ["method", "1e-4", "1e-10"],
        ["MC-ESO", "95.2", "86.2"],
        ["DE", "91.2", "89.8"],
        ["SaVOA", "71.5", "59.5"],
        ["PSO", "72.5", "48.9"],
        ["CMA-ES", "70.0", "64.3"],
    ], row_h=Inches(0.5), fontsize=13, highlight_row=1,
        cell_fills={(1, 1): PALE_GRN, (2, 2): PALE_GRN, (1, 2): PALE_RED},
        cell_text={(1, 1): GRN_TX, (2, 2): GRN_TX, (1, 2): RED})


def p6_diagnosis():
    s = slide()
    chrome(s, "Diagnosis — two weaknesses under the hood", "2 — Direction", 6)
    img(s, FIG / "p08_diagnosis/diag.emf", MARGIN, Inches(1.55), w=Inches(6.2))
    # right: two weakness notes
    x = Inches(7.2); w = Inches(5.6)
    accent_item(s, x, Inches(1.7), w,
                "Main driver = the channels",
                "Channels → isotropic search:  84 → 49%.  Restart-luck rejected.")
    accent_item(s, x, Inches(2.9), w,
                "Weakness A — ill-conditioned valleys",
                "F13 / F14 / F18  /  runs short of 1e-10", label_color=RED)
    accent_item(s, x, Inches(4.0), w,
                "Weakness B — multimodal is a facade",
                "SR 100% but PR tiny (Himmelblau 0.28, Shubert 0.06)  /  live n_elite ≈ 1",
                label_color=RED)
    # mini peak table
    table(s, x, Inches(5.35), [Inches(2.4), Inches(1.6), Inches(1.6)], [
        ["multi-global", "PR@1e-4", "SR@1e-10"],
        ["Himmelblau", "0.28", "100%"],
        ["Shubert", "0.06", "100%"],
    ], row_h=Inches(0.38), fontsize=12)


def _tint(col, f=0.82):
    """Mix a palette colour with white (f = fraction white) for a light fill."""
    h = str(col)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    m = lambda c: int(round(c + (255 - c) * f))
    return RGBColor(m(r), m(g), m(b))


def _shade(col, f=0.26):
    """Mix a palette colour toward the deck's dark slate — deepens/mutes a hue so
    bright accents (e.g. AMBER) sit in the deck's restrained red+neutral tone."""
    h = str(col)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    d = (0x1F, 0x27, 0x33)
    m = lambda c, dc: int(round(c + (dc - c) * f))
    return RGBColor(m(r, d[0]), m(g, d[1]), m(b, d[2]))


def _dir_card(s, x, w, hdr, tag, pill, pill_col, bullets, acc,
              foot, foot_fill, foot_col):
    top = Inches(2.15); h = Inches(3.95); hh = Inches(0.82)
    # rounded body; header band rounds its TOP corners to match the card and is
    # squared along the bottom where it meets the body
    softbox(s, x, top, w, h, fill=WHITE, line=BORDER, line_w=1.2)
    softbox(s, x, top, w, hh, fill=hdr)
    rect(s, x, top + hh - Inches(BOXR_IN), w, Inches(BOXR_IN), fill=hdr)
    txt(s, x + Inches(0.32), top, w - Inches(2.4), hh, tag, size=21, bold=True,
        color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # status pill on the header right
    pw = Inches(1.55)
    softbox(s, x + w - pw - Inches(0.28), top + Inches(0.21), pw, Inches(0.4),
            fill=WHITE)
    txt(s, x + w - pw - Inches(0.28), top + Inches(0.21), pw, Inches(0.4), pill,
        size=12, bold=True, color=pill_col, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    # keyword bullets
    lines = []
    for i, b in enumerate(bullets):
        lines.append({"runs": [("›  ", 17, True, acc), (b, 16.5, False, DARK)],
                      "space_before": 0 if i == 0 else 14})
    paras(s, x + Inches(0.38), top + hh + Inches(0.42), w - Inches(0.76),
          Inches(1.9), lines)
    # footer note strip
    fh = Inches(0.66)
    softbox(s, x + Inches(0.32), top + h - fh - Inches(0.3), w - Inches(0.64), fh,
            fill=foot_fill)
    txt(s, x + Inches(0.32), top + h - fh - Inches(0.3), w - Inches(0.64), fh,
        foot, size=14, bold=True, color=foot_col, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)


def p7_direction():
    s = slide()
    chrome(s, "Performance first", "2 — Direction")
    subtitle(s, "Fix the SR@1e-10 gaps before chasing multiple optima")
    _dir_card(
        s, MARGIN, Inches(5.85), RED, "A. Performance", "CHOSEN", RED,
        ["Fill deep-precision gaps on F11–F14", "SR@1e-10 never lowered"],
        RED, "Primary metric = SR@1e-10", PINK, RED_DK)
    _dir_card(
        s, Inches(6.95), Inches(5.85), GRAY, "B. Multimodality", "PARTIAL", GRAY,
        ["Find many optima in parallel", "Sequential niching — SR-safe only"],
        GRAY, "Deeper work deferred  →  §5", THEAD, GRAY)


def p8_timeline():
    s = slide()
    chrome(s, "Four changes since 5/18", "3 — Improvements")
    subtitle(s, "Cumulative — each builds on the last")
    items = [("6/12", "Informed restart", "reservoir re-ignition + basin repulsion",
              "p10_timeline/restart.emf"),
             ("6/24", "Adaptive floor", "eigenvalue-ratio anisotropy", "p10_timeline/floor.emf"),
             ("7/03", "Channel router", "per-landscape air-budget routing", "p10_timeline/router.emf"),
             ("7/03", "Route-gated best2", "2nd difference on droplet route", "p10_timeline/best2.emf")]
    cw = Inches(2.85); x0 = Inches(0.5); gap = Inches(0.28)
    line_y = Inches(2.55); ctop = Inches(2.9); ch = Inches(3.0)
    fw = Inches(2.2)
    centers = [x0 + i * (cw + gap) + cw / 2 for i in range(len(items))]
    # horizontal timeline spanning the full card row (runs past the end nodes)
    line_l = x0
    line_r = x0 + (len(items) - 1) * (cw + gap) + cw
    hline(s, line_l, line_y, line_r - line_l, color=PINK, weight=3.5)
    nd = Inches(0.17)
    for i, (dt, nm, d, fn) in enumerate(items):
        x = x0 + i * (cw + gap)
        c = centers[i]
        # date above the node
        txt(s, x, line_y - Inches(0.62), cw, Inches(0.36), dt, size=15, bold=True,
            color=RED, align=PP_ALIGN.CENTER)
        # node dot on the line
        rect(s, c - nd / 2, line_y - nd / 2, nd, nd, fill=RED, shape=MSO_SHAPE.OVAL)
        # card below
        softbox(s, x, ctop, cw, ch, fill=WHITE, line=BORDER, line_w=1.0)
        img(s, FIG / fn, x + (cw - fw) / 2, ctop + Inches(0.25), w=fw)
        paras(s, x + Inches(0.2), ctop + Inches(1.78), cw - Inches(0.4), Inches(1.1), [
            {"runs": [(nm, 16, True, DARK)], "align": PP_ALIGN.CENTER},
            {"runs": [(d, 12, False, GRAY)], "space_before": 6,
             "align": PP_ALIGN.CENTER},
        ])


def p9_waterfall():
    s = slide()
    chrome(s, "Cumulative ablation — +6.0 pt", "3 — Improvements")
    subtitle(s, "SR@1e-10 86.9 → 92.9%")
    datanote(s, "BBOB-24 / dim 2 / n = 20", Inches(1.85))
    pw = Inches(5.35)
    img(s, FIG / "p11_waterfall/sr1e4.emf", Inches(0.75), Inches(1.85), w=pw)
    img(s, FIG / "p11_waterfall/sr1e10.emf", Inches(6.6), Inches(1.85), w=pw)
    # this cycle's contribution
    softbox(s, MARGIN, Inches(6.12), Inches(12.23), Inches(0.6), fill=BG,
            line=RED, line_w=1.3)
    paras(s, Inches(0.85), Inches(6.12), Inches(11.6), Inches(0.6), [
        {"runs": [("Now strong at both ends —  ", 14, True, RED_DK),
                  ("coarse discovery kept (1e-4 saturated), deep precision gained "
                   "(1e-10 +6, past DE)", 13.5, False, DARK)]},
    ], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, MARGIN, Inches(6.82), Inches(12.23), Inches(0.3),
        "Sanity: abl0 86.9% ≈ the 86.2% measured on 5/18 (n=30)",
        size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def _hexc(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def p10_methods():
    s = slide()
    chrome(s, "Standing today — first at deep precision", "4 — Comparison")
    subtitle(s, "SR@1e-10 92.9% — DE overtaken")
    datanote(s, "each cell: SR@1e-10 (left) / evals-to-success (right) / "
                "BBOB-24 / dim 2 / n = 20", Inches(2.2))
    # combined table (left): rows = methods, cols = categories, split L/R.
    # No on-slide legend (explained verbally); table is centred vertically.
    iw = Inches(8.3)
    tx = Inches(0.1)
    img(s, FIG / "p31_category/catsplit.emf", tx, Inches(2.2), w=iw)
    # ── divider between the table and the Wilcoxon panel ─────────────────────
    rect(s, Inches(8.6), Inches(2.35), Emu(13000), Inches(4.05), fill=RULEC)
    # ── Wilcoxon panel (right, narrowed): W / T / L vs each baseline ──────────
    WIN = GREEN; TIE = RGBColor(0xD3, 0xD8, 0xDE); LOSS = RED
    wx = Inches(8.9); ww = Inches(4.15)
    txt(s, wx, Inches(2.5), ww, Inches(0.32), "Wilcoxon — vs each baseline",
        size=12.5, bold=True, color=DARK)
    lx2 = wx
    for lab, col in [("win", WIN), ("tie", TIE), ("loss", LOSS)]:
        rect(s, lx2, Inches(2.95), Inches(0.16), Inches(0.16), fill=col)
        txt(s, lx2 + Inches(0.22), Inches(2.89), Inches(0.7), Inches(0.3), lab,
            size=10, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        lx2 = lx2 + Inches(0.78)
    # the unified comparison set, hardest-to-beat first (from wilcoxon.csv, α=0.05)
    wtl = [("DE", 2, 21, 1), ("IPOP-CMA-ES", 6, 18, 0),
           ("NM-Restart", 10, 13, 1), ("SaVOA", 11, 12, 1), ("PSO", 18, 6, 0)]
    lab_w = Inches(1.35); bar_w = Inches(2.6)
    unit = int(bar_w) / 24
    bh = Inches(0.3); rh = Inches(0.44); gap = Inches(0.16)
    y = Inches(3.55)
    for name, wv, tv, lv in wtl:
        txt(s, wx, y, lab_w, rh, name, size=10.5, bold=True, color=DARK,
            anchor=MSO_ANCHOR.MIDDLE)
        bx = wx + lab_w
        for cnt, col, tc in [(wv, WIN, WHITE), (tv, TIE, DARK), (lv, LOSS, WHITE)]:
            if cnt > 0:
                seg = Emu(int(round(cnt * unit)))
                rect(s, bx, y + Inches(0.04), seg, bh, fill=col)
                txt(s, bx, y + Inches(0.04), seg, bh, str(cnt), size=10,
                    bold=True, color=tc, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
                bx = Emu(int(bx) + int(seg))
        y = Emu(int(y) + int(rh) + int(gap))


def p12_restart():
    s = slide()
    chrome(s, "① Informed restart", "3 — Improvements")
    subtitle(s, "A stalled sub-population re-seeds — real run on F04")
    pw = Inches(4.5)
    img(s, FIG / "p12_restart/before.emf", Inches(1.4), Inches(1.7), w=pw)
    img(s, FIG / "p12_restart/after.emf", Inches(7.4), Inches(1.7), w=pw)


def p12b_restart_math():
    s = slide()
    chrome(s, "① Informed restart — how it works", "3 — Improvements")
    subtitle(s, "Re-seed from what the search already found, not blindly")
    # ── REPELLED (left) + RE-IGNITE (right); both run — a split, not a choice.
    #    Each mode is ONE card (colour header = name, body = its formula). A big
    #    "+" and the firing condition sit dead-centre between them ──────────────
    cx = SW / 2
    lcx_in, rcx_in = 3.55, 9.78          # left = REPELLED · right = RE-IGNITE
    lcx, rcx = Inches(lcx_in), Inches(rcx_in)
    # merged mode card: colour header (name) directly above its formula body
    ct, cardw, cardh, hh = Inches(1.92), Inches(4.3), Inches(1.06), Inches(0.42)
    def _mode(bx_in, header, col, eqname):
        cl = Inches(bx_in) - cardw / 2
        softbox(s, cl, ct, cardw, cardh, fill=WHITE, line=col)
        softbox(s, cl, ct, cardw, hh, fill=col)             # rounded-top header
        rect(s, cl, ct + hh - Inches(BOXR_IN), cardw, Inches(BOXR_IN), fill=col)
        txt(s, cl, ct, cardw, hh, header, size=13, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _eq_center(s, eqname, bx_in, 2.66, folder="p12b_restart_math")
    _mode(lcx_in, "REPELLED", BLUE, "repel")
    _mode(rcx_in, "RE-IGNITE", RED, "reig")
    # dead-centre "+" (both run) and, just below it, the firing condition
    txt(s, cx - Inches(0.9), Inches(3.34), Inches(1.8), Inches(1.0), "+",
        size=72, bold=True, color=GRAY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    paras(s, cx - Inches(1.9), Inches(4.32), Inches(3.8), Inches(0.92), [
        {"align": PP_ALIGN.CENTER,
         "runs": [("no improvement 300 evals", 11.5, True, DARK)]},
        {"align": PP_ALIGN.CENTER, "space_before": 4,
         "runs": [("&", 11.5, True, DARK)]},
        {"align": PP_ALIGN.CENTER, "space_before": 4,
         "runs": [("f_best/|f_init| > 1e-8", 11.5, True, DARK)]},
    ], anchor=MSO_ANCHOR.MIDDLE)
    # the two mechanism panels, one per mode, each with a grey caption line
    pw, capw = Inches(2.6), Inches(4.2)
    img(s, FIG / "p12b_restart_math/panel_repel.png", lcx - pw / 2,
        Inches(3.14), w=pw)
    img(s, FIG / "p12b_restart_math/panel_reig.png", rcx - pw / 2,
        Inches(3.14), w=pw)
    for bx, cap in [(lcx, "uniform, but repelled from remembered basins"),
                    (rcx, "re-light around archived elites — good spots found")]:
        txt(s, bx - capw / 2, Inches(5.86), capw, Inches(0.32), cap,
            size=13, color=GRAY, align=PP_ALIGN.CENTER)


def why_col(s, x, y, w, tag, bullets):
    """A 'Why it works' / 'Result' bullet column with a red tick header."""
    rect(s, x, y + Inches(0.03), Inches(0.08), Inches(0.5), fill=RED)
    txt(s, x + Inches(0.24), y, w, Inches(0.5), tag, size=18, bold=True, color=DARK)
    paras(s, x + Inches(0.24), y + Inches(0.7), w, Inches(3.0),
          [{"runs": [("›  ", 16, True, RED), (b, 15.5, False, DARK)],
            "space_before": 0 if i == 0 else 13} for i, b in enumerate(bullets)])


def p13b_restart_why():
    import numpy as np
    d = np.load(FIG / "restart_conv.npz", allow_pickle=True)
    conv_result_slide(
        "① Informed restart — result",
        "Blind restart stays stuck — informed restart escapes",
        "restart_conv.npz", "p13_restart_result",
        [("a", str(d["a_name"])), ("b", str(d["b_name"]))])


def p13c_restart_bar():
    s = slide()
    chrome(s, "① Informed restart — per-function SR", "3 — Improvements")
    subtitle(s, "Net ±0 on BBOB-24 — the payoff is multimodal escape")
    datanote(s, "change-ablation / BBOB-24 / dim 2 / n = 20", Inches(2.05))
    iw = Inches(9.8)
    img(s, FIG / "p14_restart_bar/sr.emf", (SW - iw) / 2, Inches(2.05), w=iw)


def p15_floor():
    s = slide()
    chrome(s, "② Adaptive anisotropy floor",
           "3 — Improvements")
    subtitle(s, "Close-contact noise adapts to the landscape shape")
    pw = Inches(4.75)
    img(s, FIG / "p15_floor/illcond.emf", Inches(1.25), Inches(1.85), w=pw)
    img(s, FIG / "p15_floor/rugged.emf", Inches(7.3), Inches(1.85), w=pw)


def _eq_center(s, name, cx, ycen, folder="p16_floor_math"):
    """Place an equation image at native size, centred on (cx, ycen) inches, so
    every equation keeps the same font regardless of its crop height."""
    pic = s.shapes.add_picture(str(FIG / f"{folder}/{name}.png"),
                               Inches(cx), Inches(ycen))
    pic.left = int(Inches(cx) - pic.width / 2)
    pic.top = int(Inches(ycen) - pic.height / 2)
    return pic


def p15c_floor_math():
    s = slide()
    chrome(s, "② Adaptive anisotropy floor — the math", "3 — Improvements")
    subtitle(s, "Only the floor φ is new — the pipeline is unchanged")
    PINKBG = RGBColor(0xFC, 0xF2, 0xF0)   # the one intentional emphasis fill

    def _softrect(s_, l_, t_, w_, h_, **k):    # deck-wide corner radius (softbox)
        return softbox(s_, l_, t_, w_, h_, **k)

    # ── Zone 1: the unchanged 3-step pipeline ────────────────────────────
    txt(s, MARGIN, Inches(1.74), Inches(12.2), Inches(0.3),
        "Close-contact channel — same flow as before", size=12, italic=True,
        color=GRAY)
    cards = [
        (0.55, 3.6, "EIGEN-DECOMPOSE", "eq1",
         "population spread → axes V, variances λ", False),
        (4.87, 3.6, "FLOOR THE AXES", "eq4",
         "lift any collapsed axis up to the floor φ", True),
        (9.19, 3.6, "SAMPLE CHILD", "eq5",
         "draw from the floored ellipse (step σ_i, noise z)", False),
    ]
    cy, ch = 2.1, 1.42
    for x, w, title, eq, cap, hl in cards:
        _softrect(s, Inches(x), Inches(cy), Inches(w), Inches(ch), fill=BG,
                  line=(RED if hl else BORDER))
        txt(s, Inches(x), Inches(cy + 0.14), Inches(w), Inches(0.3),
            title, size=12.5, bold=True, color=(RED if hl else DARK),
            align=PP_ALIGN.CENTER)
        _eq_center(s, eq, x + w / 2, cy + 0.66)
        txt(s, Inches(x + 0.15), Inches(cy + 1.05), Inches(w - 0.3), Inches(0.34),
            cap, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    for ax in (4.32, 8.64):     # block arrows between the cards
        rect(s, Inches(ax), Inches(cy + 0.5), Inches(0.4), Inches(0.34),
             fill=RGBColor(0xC7, 0xCD, 0xD4), shape=MSO_SHAPE.RIGHT_ARROW)
    # ── Zone 2: what changed — φ is now adaptive ─────────────────────────
    zy = 3.82
    _softrect(s, MARGIN, Inches(zy), Inches(12.23), Inches(3.0), fill=PINKBG,
              line=RGBColor(0xE6, 0xC7, 0xC0), line_w=1.0)
    # red down-arrow linking the floor step into this zone
    rect(s, Inches(6.47), Inches(cy + ch - 0.02), Inches(0.34), Inches(0.38),
         fill=RED, shape=MSO_SHAPE.DOWN_ARROW)
    tx = htick(s, Inches(0.78), Inches(zy + 0.18))
    txt(s, tx, Inches(zy + 0.2), Inches(6.4), Inches(0.42),
        "The floor φ is now adaptive", size=16, bold=True, color=DARK)
    # left: derivation r → ρ → φ
    _eq_center(s, "eq2", 3.9, zy + 0.98)
    txt(s, Inches(0.85), Inches(zy + 1.34), Inches(6.1), Inches(0.34),
        "EMA = exponential moving average — smooths out transient blips",
        size=10.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.74), Inches(zy + 1.68), Inches(0.32), Inches(0.3),
         fill=RGBColor(0xC7, 0xCD, 0xD4), shape=MSO_SHAPE.DOWN_ARROW)
    _eq_center(s, "eq3", 3.9, zy + 2.3)
    # right: illustration of the adaptive floor (transparent, two cases)
    img(s, FIG / "p16_floor_math/illus.png", Inches(6.9), Inches(zy + 0.72),
        w=Inches(5.2))


def conv_result_slide(title, subtitle_text, npz, folder, rows):
    """Improved-seed result slide: per function a small 2-D map + 3-D landscape
    with the function label above them, and a large convergence panel on the
    right. Shared by the floor (②) and router (③) result slides.
    rows = [(tag, "Function name"), ...] (two rows)."""
    import numpy as np
    s = slide()
    chrome(s, title, "3 — Improvements")
    d = np.load(FIG / npz, allow_pickle=True)
    # no subtitle — the freed space becomes vertical margin above / between /
    # below the two convergence rows
    xm, wm = Inches(1.4), Inches(1.9)      # 2-D map (small square)
    xd, wd = Inches(3.45), Inches(2.0)     # 3-D landscape (small)
    xc, wc = Inches(5.6), Inches(6.1)      # convergence (large, wide)
    for (tag, name), y in zip(rows, [Inches(1.55), Inches(4.42)]):
        txt(s, xm, y, Inches(4.05), Inches(0.46),
            f"{name}  /  seed {int(d[f'{tag}_seed'])}", size=12, bold=True,
            color=DARK, align=PP_ALIGN.CENTER)
        img(s, FIG / f"{folder}/{tag}_map.emf", xm, y + Inches(0.4), w=wm)
        img(s, FIG / f"{folder}/{tag}_surf.emf", xd, y + Inches(0.46), w=wd)
        img(s, FIG / f"{folder}/{tag}_conv.emf", xc, y + Inches(0.1), w=wc)


def p15b_floor_why():
    conv_result_slide(
        "② Adaptive anisotropy floor — result",
        "Kills the rare runs that stall just above 1e-10",
        "floor_conv.npz", "p17_floor_result",
        [("f10", "F10-EllipsoidalRot"), ("f19", "F19-GriewankRosenbrock")])


def p17b_floor_bar():
    s = slide()
    chrome(s, "② Adaptive anisotropy floor — per-function SR", "3 — Improvements")
    subtitle(s, "Net +3.1 pt — the largest single step")
    datanote(s, "change-ablation / BBOB-24 / dim 2 / n = 20 / "
                "functions whose SR@1e-10 moved", Inches(2.05))
    iw = Inches(9.8)
    img(s, FIG / "p18_floor_result_bar/sr.emf", (SW - iw) / 2, Inches(2.05), w=iw)


def p18c_router_conv():
    import numpy as np
    d = np.load(FIG / "router_conv.npz", allow_pickle=True)
    conv_result_slide(
        "③ Channel router — result",
        "Routed runs converge; unrouted ones stall",
        "router_conv.npz", "p21_router_conv",
        [("a", str(d["a_name"])), ("b", str(d["b_name"]))])


def detail_slide(page, title, changed, why, result, panel_fn,
                 section="3 — Improvements"):
    s = slide()
    chrome(s, title, section, page)
    # left concept panel
    panel_fn(s, MARGIN, Inches(1.55), Inches(5.6), Inches(4.9))
    # right text blocks — body may be a str (one line) or a list (keyword bullets)
    x = Inches(6.55); w = Inches(6.25)
    y = Inches(1.7)
    for tag, body, col in [("What changed", changed, RED),
                           ("Why it works", why, RED),
                           ("Result", result, RED)]:
        lines = [{"runs": [(tag, 13, True, col)]}]
        if isinstance(body, (list, tuple)):
            for b in body:
                lines.append({"runs": [("›  ", 15, True, col), (b, 15, False, DARK)],
                              "space_before": 4})
        else:
            lines.append({"runs": [(body, 15, False, DARK)], "space_before": 4,
                          "spacing": 1.05})
        paras(s, x, y, w, Inches(1.5), lines)
        y = Emu(int(y) + int(Inches(1.6)))


def panel_ir(s, l, t, w, h):
    rect(s, l, t, w, h, fill=BG, line=BORDER, line_w=1.0)
    half = Emu(int(w) // 2)
    txt(s, l, t + Inches(0.15), half, Inches(0.4), "before", size=13, bold=True,
        color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, l + half, t + Inches(0.15), half, Inches(0.4), "after", size=13, bold=True,
        color=RED, align=PP_ALIGN.CENTER)
    import random
    random.seed(3)
    # before: scattered uniform dots
    bl = l + Inches(0.4); bt = t + Inches(0.75); bw = Emu(int(half) - int(Inches(0.8)))
    bh = Emu(int(h) - int(Inches(1.1)))
    for _ in range(22):
        dx = int(bl) + int(random.random() * int(bw))
        dy = int(bt) + int(random.random() * int(bh))
        rect(s, Emu(dx), Emu(dy), Inches(0.09), Inches(0.09), fill=GRAY,
             shape=MSO_SHAPE.OVAL)
    # after: cluster near reservoir + repel ring
    al = l + half + Inches(0.4)
    cx = int(al) + int(bw) // 2; cy = int(bt) + int(bh) // 2
    rect(s, Emu(cx - int(Inches(0.65))), Emu(cy - int(Inches(0.65))),
         Inches(1.3), Inches(1.3), fill=None, line=RED, line_w=1.5,
         shape=MSO_SHAPE.OVAL)
    for _ in range(14):
        ang = random.random() * 6.28; rr = random.random() * int(Inches(0.5))
        import math
        dx = cx + int(rr * math.cos(ang)); dy = cy + int(rr * math.sin(ang))
        rect(s, Emu(dx), Emu(dy), Inches(0.09), Inches(0.09), fill=RED,
             shape=MSO_SHAPE.OVAL)
    txt(s, Emu(cx - int(Inches(0.9))), Emu(cy + int(Inches(0.7))), Inches(1.8),
        Inches(0.3), "avoid abandoned basin", size=10.5, color=GRAY,
        align=PP_ALIGN.CENTER)


def panel_floor(s, l, t, w, h):
    rect(s, l, t, w, h, fill=BG, line=BORDER, line_w=1.0)
    paras(s, l + Inches(0.3), t + Inches(0.35), w - Inches(0.6), Inches(4.2), [
        {"runs": [("Anisotropy floor set by eigenvalue ratio", 15, True, DARK)]},
        {"runs": [("λmax / λmin  (scale- & shift-invariant)", 12.5, False, GRAY)],
         "space_before": 3},
        {"runs": [("ill-conditioned  1e5–1e7", 14, True, RED)], "space_before": 20},
        {"runs": [("→ lower the floor, release anisotropy along the valley",
                   13, False, DARK)], "space_before": 3},
        {"runs": [("rugged / multimodal  3–600", 14, True, BLUE)], "space_before": 16},
        {"runs": [("→ keep the floor high, clamp false anisotropy",
                   13, False, DARK)], "space_before": 3},
    ])


def panel_router(s, l, t, w, h):
    rect(s, l, t, w, h, fill=BG, line=BORDER, line_w=1.0)
    txt(s, l, t + Inches(0.2), w, Inches(0.4), "covariance signals → one route",
        size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    routes = [("cond > 3", "DROPLET", "ill-cond valley", RED),
              ("algA & mgap high", "CLOSE", "separable", BLUE),
              ("otherwise", "KEEP-AIR", "multimodal = base", GREEN)]
    y = t + Inches(0.8)
    for cond, route, note, col in routes:
        rect(s, l + Inches(0.35), y, Inches(2.3), Inches(0.85), fill=WHITE,
             line=BORDER, line_w=1.0, shape=MSO_SHAPE.RECTANGLE)
        txt(s, l + Inches(0.45), y, Inches(2.1), Inches(0.85), cond, size=12.5,
            color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        # arrow
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l + Inches(2.75), y + Inches(0.28),
                                Inches(0.55), Inches(0.3))
        ar.fill.solid(); ar.fill.fore_color.rgb = RULEC; ar.line.fill.background()
        ar.shadow.inherit = False
        rect(s, l + Inches(3.4), y, Inches(1.7), Inches(0.85), fill=col,
             shape=MSO_SHAPE.RECTANGLE)
        paras(s, l + Inches(3.4), y + Inches(0.12), Inches(1.7), Inches(0.65), [
            {"runs": [(route, 14, True, WHITE)], "align": PP_ALIGN.CENTER},
            {"runs": [(note, 9.5, False, WHITE)], "align": PP_ALIGN.CENTER,
             "space_before": 1},
        ])
        y = Emu(int(y) + int(Inches(1.25)))


def panel_best2(s, l, t, w, h):
    rect(s, l, t, w, h, fill=BG, line=BORDER, line_w=1.0)
    cx = int(l) + int(w) // 2
    txt(s, l, t + Inches(0.2), w, Inches(0.4),
        "droplet difference structure", size=13, bold=True, color=DARK,
        align=PP_ALIGN.CENTER)
    # single vs double difference vectors
    paras(s, l + Inches(0.4), t + Inches(1.0), w - Inches(0.8), Inches(1.2), [
        {"runs": [("single (base):  ", 15, False, GRAY),
                  ("xp + F(x_strain − xp)", 15, True, GRAY)]},
        {"runs": [("route-gated:  ", 15, False, DARK),
                  ("… + F(x_c − x_d)", 16, True, RED)], "space_before": 22},
    ])
    paras(s, l + Inches(0.4), t + Inches(3.0), w - Inches(0.8), Inches(1.4), [
        {"runs": [("2nd diff = donor diversity  →  escape wrong basin", 13, False, GRAY)]},
        {"runs": [("droplet route only / off-route bit-identical", 13, False, GRAY)],
         "space_before": 6},
    ])


def panel_niching(s, l, t, w, h):
    """Two-regime schematic for sequential niching: a run behaves exactly like
    base until it drills the current basin out (σ-exhaustion), then a repelled
    restart seeds the next unfound optimum. Width-responsive so it can share the
    slide with the PR/SR charts on the right."""
    GOLD = RGBColor(0xF2, 0xC1, 0x4E)
    W = int(w) / int(Inches(1))                  # panel width in inches
    cm = 0.26                                     # card inset
    cw = W - 2 * cm                               # card width
    scx = cm + cw * 0.79                          # schematic column centre-x
    txtw = cw * 0.54                              # left text-block width

    def at(dx, dy):                              # (inches from panel origin) → EMU
        return Emu(int(l) + int(Inches(dx))), Emu(int(t) + int(Inches(dy)))

    def star(dx, dy, sz=0.32, fill=GOLD):
        x, y = at(dx - sz / 2, dy - sz / 2)
        rect(s, x, y, Inches(sz), Inches(sz), fill=fill,
             shape=MSO_SHAPE.STAR_5_POINT)

    def dot(dx, dy, col=GRAY, sz=0.1):
        x, y = at(dx - sz / 2, dy - sz / 2)
        rect(s, x, y, Inches(sz), Inches(sz), fill=col, shape=MSO_SHAPE.OVAL)

    rect(s, l, t, w, h, fill=BG, line=BORDER, line_w=1.0)
    txt(s, l, t + Inches(0.16), w, Inches(0.34), "One run / two regimes",
        size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    ch = 1.42                                     # card height
    # ── regime ① — base (grey): drill the current basin, unchanged ──────────
    y1 = 0.60
    bx, by = at(cm, y1)
    softbox(s, bx, by, Inches(cw), Inches(ch), fill=WHITE, line=BORDER, line_w=1.0)
    paras(s, *at(cm + 0.22, y1 + 0.20), Inches(txtw), Inches(1.05), [
        {"runs": [("① Base regime", 13, True, GRAY)]},
        {"runs": [("drill one basin to 1e-10", 12, False, DARK)],
         "space_before": 4},
        {"runs": [("— identical to base", 11, False, GRAY)], "space_before": 2},
    ])
    c1 = y1 + ch / 2
    for dx, dy in [(scx - 0.2, c1 - 0.34), (scx + 0.24, c1 - 0.12),
                   (scx - 0.26, c1 + 0.30), (scx + 0.16, c1 + 0.28)]:
        dot(dx, dy)
    star(scx, c1)

    # ── the gate between regimes ────────────────────────────────────────────
    gy = y1 + ch + 0.06
    ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, *at(cm + 0.30, gy),
                            Inches(0.30), Inches(0.40))
    ar.fill.solid(); ar.fill.fore_color.rgb = RED; ar.line.fill.background()
    ar.shadow.inherit = False
    txt(s, *at(cm + 0.74, gy - 0.02), Inches(cw - 0.8), Inches(0.42),
        "σ-exhausted  (scale-invariant)", size=11, bold=True, color=RED,
        anchor=MSO_ANCHOR.MIDDLE)

    # ── regime ② — niching (red): repelled restart → next optimum ───────────
    y2 = gy + 0.60
    bx, by = at(cm, y2)
    softbox(s, bx, by, Inches(cw), Inches(ch), fill=WHITE, line=RED, line_w=1.3)
    paras(s, *at(cm + 0.22, y2 + 0.20), Inches(txtw), Inches(1.05), [
        {"runs": [("② Niching regime", 13, True, RED)]},
        {"runs": [("repel restart to the", 12, False, DARK)], "space_before": 4},
        {"runs": [("next unfound optimum", 12, False, DARK)]},
    ])
    c2 = y2 + ch / 2
    star(scx - 0.42, c2)                          # optimum just captured
    aj = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, *at(scx - 0.22, c2 - 0.13),
                            Inches(0.48), Inches(0.26))
    aj.fill.solid(); aj.fill.fore_color.rgb = RED; aj.line.fill.background()
    aj.shadow.inherit = False
    star(scx + 0.42, c2)                          # next optimum, freshly seeded


def p16_evals():
    s = slide()
    chrome(s, "Ill-conditioned valleys converge ~2× faster",
           "3 — Improvements")
    subtitle(s, "Same SR, about half the evaluations")
    datanote(s, "evals-to-success / same ablation ladder / BBOB-24 / dim 2 / "
                "n = 20", Inches(2.1))
    iw = Inches(12.0)
    img(s, FIG / "p30_evals/evals.emf", (SW - iw) / 2, Inches(2.1), w=iw)


def p18_multimodal_diag():
    s = slide()
    chrome(s, "Multimodality — SR hides the gap", "5 — Multimodality")
    subtitle(s, "SR is 100%, yet few optima are found (low PR)")
    funcs = [
        ("himmelblau", "C01 Himmelblau", "4", "0.28"),
        ("sixhump", "C02 Six-hump", "2", "0.60"),
        ("shubert", "C03 Shubert", "18", "0.06"),
    ]
    # three larger examples shown as a clean gallery (no bottom panel)
    cw, gap, mw = 3.86, 0.28, 2.55
    x0 = MARGIN + (Inches(12.23) - Inches(3 * cw + 2 * gap)) / 2
    cy = Inches(2.0)
    for i, (tag, name, kopt, pr) in enumerate(funcs):
        cx = int(x0) + int(Inches(i * (cw + gap) + cw / 2))
        img(s, FIG / f"p32_multimodal/shape_{tag}.emf",
            Emu(cx - int(Inches(mw / 2))), cy, w=Inches(mw))
        yb = int(cy) + int(Inches(mw + 0.18))
        cl = Emu(cx - int(Inches(cw / 2)))
        txt(s, cl, Emu(yb), Inches(cw), Inches(0.4), name,
            size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        txt(s, cl, Emu(yb + int(Inches(0.44))), Inches(cw), Inches(0.32),
            f"K = {kopt} optima", size=12.5, color=GRAY, align=PP_ALIGN.CENTER)
        # the facade: SR full (green) vs PR low (red)
        paras(s, cl, Emu(yb + int(Inches(0.92))), Inches(cw), Inches(0.4),
              [{"align": PP_ALIGN.CENTER,
                "runs": [("SR@1e-10  ", 12, False, GRAY),
                         ("100%", 17, True, GRN_DK)]}])
        paras(s, cl, Emu(yb + int(Inches(1.42))), Inches(cw), Inches(0.4),
              [{"align": PP_ALIGN.CENTER,
                "runs": [("PR@1e-4  ", 12, False, GRAY),
                         (pr, 17, True, RED)]}])


def p19_pr_vs_sr():
    s = slide()
    chrome(s, "Sequential niching — more optima, SR untouched",
           "5 — Multimodality", 19)
    subtitle(s, "One run captures optima one at a time — SR@1e-10 stays 100%")
    # left: the mechanism on a real landscape (Himmelblau, 4 optima in sequence)
    fw = Inches(4.2)
    fx = MARGIN + (Inches(4.6) - fw) / 2
    img(s, FIG / "p34_niching/sequence.emf", fx, Inches(1.95), w=fw)
    # two-regime rule, spelled out under the figure
    paras(s, MARGIN, Inches(6.0), Inches(4.6), Inches(1.0), [
        {"runs": [("① drill a basin to 1e-10", 13, True, DARK),
                  ("  — like base", 12, False, GRAY)]},
        {"runs": [("② σ-exhausted → repelled restart", 13, True, RED),
                  ("  to the next optimum", 12, False, GRAY)],
         "space_before": 8},
    ])
    # right: PR-by-function bar chart (base vs MC-ESO), SR stays 100% ────────
    img(s, FIG / "p33_pr_vs_sr/pr.emf", Inches(5.15), Inches(2.15), w=Inches(7.55))


def p20_multimodal_remain():
    s = slide()
    chrome(s, "Multimodality — still only partial", "5 — Multimodality", 20)
    accent_item(s, MARGIN, Inches(1.7), Inches(12), "Full multi-solution not reached",
                "MMO success @1e-4 still 0% on Himmelblau / Shubert — not every optimum captured",
                label_color=RED)
    accent_item(s, MARGIN, Inches(2.9), Inches(12), "The dilemma",
                "Parallel basin drilling collapses SR@1e-10  /  deep precision vs full coverage")
    accent_item(s, MARGIN, Inches(4.1), Inches(12), "Deliberate scope",
                "Performance axis prioritized  /  multimodal advanced only where SR-free")
    rect(s, MARGIN, Inches(5.5), Inches(12.23), Inches(0.9), fill=BG, line=BORDER,
         line_w=1.0)
    txt(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(0.9),
        "Shubert PR@1e-4 tripled (0.06 → 0.17) — real progress, absolute level still low",
        size=14, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)


def p21_summary():
    s = slide()
    chrome(s, "Summary", "6 — Wrap-up", 21)
    subtitle(s, "The whole cycle in three lines")
    # TOP: Idea / Mechanism / Result side by side — accent-bar heading + its
    # contents as chevron bullets (one item per bullet)
    sections = [
        ("Idea", [[("Uniform channel mix  →", 14, False, DARK)],
                  [("landscape-aware routing", 14, False, DARK)]]),
        ("Mechanism", [[("informed restart", 14, False, DARK)],
                       [("adaptive floor", 14, False, DARK)],
                       [("channel router", 14, False, DARK)],
                       [("route-gated best2", 14, False, DARK)]]),
        ("Result", [[("SR@1e-10  ", 14, False, DARK), ("86.9 → 92.9%", 14, True, RED)],
                    [("top of 10 methods", 14, False, DARK)],
                    [("valleys 2× faster", 14, False, DARK)],
                    [("multi-solution ↑, SR kept", 14, False, DARK)]]),
    ]
    hy = 1.66
    for (tag, bullets), x in zip(sections, (0.55, 4.72, 8.89)):
        tx = htick(s, Inches(x), Inches(hy + 0.02))
        txt(s, tx, Inches(hy), Inches(3.6), Inches(0.36), tag, size=17, bold=True,
            color=RED, anchor=MSO_ANCHOR.MIDDLE)
        lines = [{"runs": [("›  ", 13, True, RED)] + b,
                  "space_before": 0 if i == 0 else 4} for i, b in enumerate(bullets)]
        paras(s, tx, Inches(hy + 0.42), Inches(3.5), Inches(1.9), lines)
    # BOTTOM: the two headline results — SR@1e-10 lift (left) + standing vs the
    # whole field by category (right, the p31 table)
    img(s, FIG / "p11_waterfall/sr1e10.emf", Inches(0.75), Inches(3.66),
        w=Inches(4.0))
    img(s, FIG / "p31_category/catsplit.emf", Inches(5.25), Inches(3.7),
        w=Inches(7.0))


def p22_future():
    s = slide()
    chrome(s, "Remaining work", "6 — Wrap-up", 22)
    subtitle(s, "Three directions")
    centers = [2.42, 6.667, 10.91]
    # (accent colour, title, (kind, path), figure width, caption) per card
    # deepened/muted palette so the bands sit in the deck's restrained tone
    cards = [
        (_shade(AMBER), "Generalize the constants", ("png", "p37_future/params.png"),
         3.3, "thresholds may be BBOB-tuned — make them scale-derived"),
        (_shade(TEAL), "Robustness", ("emf", "p37_future/robustness.emf"),
         2.7, "partial multi-solution  /  robust to mild noise, severe open"),
        (_shade(BLUE), "Higher dimensions", ("png", "p37_future/dim.png"),
         3.3, "at 3D the CMA-ES restart family overtakes MC-ESO at 1e-10"),
    ]
    # p7-style cards: a coloured header band (white title) · figure body ·
    # tinted footer note strip
    cwid, by, bh, hh, fh = 3.86, 1.86, 4.88, 0.62, 0.86
    for cx, (col, title, (kind, fp), fw, cap) in zip(centers, cards):
        bx = Inches(cx - cwid / 2)
        softbox(s, bx, Inches(by), Inches(cwid), Inches(bh), fill=WHITE,
                line=col, line_w=1.2)
        # coloured header band, rounded top to match the card, squared bottom
        softbox(s, bx, Inches(by), Inches(cwid), Inches(hh), fill=col)
        rect(s, bx, Inches(by + hh - BOXR_IN), Inches(cwid), Inches(BOXR_IN),
             fill=col)
        txt(s, bx, Inches(by), Inches(cwid), Inches(hh), title, size=15,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # figure centred in the body between header and footer
        fig_h = fw * (2.5 / 3.5) if kind == "png" else fw
        body_top, body_bot = by + hh + 0.12, by + bh - fh - 0.34
        fy = body_top + (body_bot - body_top - fig_h) / 2
        img(s, FIG / fp, Inches(cx - fw / 2), Inches(fy), w=Inches(fw))
        # footer note strip (light tint of the card colour)
        fsy = Inches(by + bh - fh - 0.22)
        softbox(s, bx + Inches(0.24), fsy, Inches(cwid - 0.48), Inches(fh),
                fill=_tint(col))
        txt(s, bx + Inches(0.24), fsy, Inches(cwid - 0.48), Inches(fh), cap,
            size=12, bold=True, color=col, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)


def p18_router():
    s = slide()
    chrome(s, "③ Channel router", "3 — Improvements")
    subtitle(s, "Each run locks one route, read off its covariance shape")
    gap_in, m_in = 0.34, 0.55                  # column gap / left margin (inches)

    # ── per-route columns: name / target-shape landscape (left) / decision
    #    (right) / budget chart below. DROPLET/CLOSE carry a full definition, so
    #    they get wide cards; KEEP-AIR is just "= base" (no equation), so its
    #    card is trimmed and the freed width goes to the other two.
    routes = [
        ("droplet", "DROPLET", RED, 4.36, True),   # DROPLET / CLOSE share a width
        ("close", "CLOSE", BLUE, 4.36, True),
        ("keepair", "KEEP-AIR", GRN_DK, 2.82, False),
    ]
    # lay the cards out left→right with their individual widths
    x0, cards = m_in, []
    for tag, name, col, bw_card, has_eq in routes:
        cx = x0 + bw_card / 2
        cards.append((tag, name, col, bw_card, has_eq, x0, cx))
        x0 += bw_card + gap_in
    # full-height thin rules separate the columns, drawn in each gap
    for _, _, _, bw_card, _, cx0, _ in cards[:-1]:
        dx = Inches(cx0 + bw_card + gap_in / 2)
        connector(s, dx, Inches(2.0), dx, Inches(6.92), color=BORDER,
                  weight=1.2, arrow=False)
    # Each route's identity — name, the target landscape, and (for the two
    # routed channels) the committing condition — sits inside ONE route-coloured
    # card; the budget chart below is the MAIN figure of the slide.
    box_h, box_y = 1.66, 2.0
    tw = 0.95                                   # landscape thumbnail (left)
    bw = 2.82                                   # ratio chart = KEEP-AIR card width
    for tag, name, col, box_w, has_eq, bx, cx in cards:
        softbox(s, Inches(bx), Inches(box_y), Inches(box_w), Inches(box_h),
                fill=BG, line=col, line_w=1.4)
        # route name across the top of the card
        txt(s, Inches(bx), Inches(box_y + 0.12), Inches(box_w), Inches(0.30),
            name, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        if has_eq:
            # thumbnail on the left, the definition centred in the space to its
            # right — but clamped so a wide equation never overlaps the thumbnail
            img(s, FIG / f"p20_router/shape_{tag}.emf", Inches(bx + 0.24),
                Inches(box_y + 0.46), w=Inches(tw))
            pic = s.shapes.add_picture(
                str(FIG / f"p20_router/dec_{tag}.png"), Inches(0), Inches(0))
            reg_l = int(Inches(bx + 0.24 + tw + 0.15))
            reg_r = int(Inches(bx + box_w - 0.15))
            pic.left = max(reg_l, (reg_l + reg_r) // 2 - int(pic.width) // 2)
            pic.top = int(Inches(box_y + 0.95) - pic.height / 2)
        else:
            # KEEP-AIR: just the shape, centred (no equation)
            img(s, FIG / f"p20_router/shape_{tag}.emf", Inches(cx - tw / 2),
                Inches(box_y + 0.46), w=Inches(tw))
        # main: the child-count budget vs σ for this route
        img(s, FIG / f"p20_router/budget_{tag}.emf", Inches(cx - bw / 2),
            Inches(3.72), w=Inches(bw))


def p18d_router_apply():
    s = slide()
    chrome(s, "③ Per-landscape channel router", "3 — Improvements")
    subtitle(s, "Re-allocate the channel child-count budget — the airborne share migrates to the routed channel as σ shrinks")
    # the child-count budget vs σ — the message of the slide, given the room
    iw = Inches(9.42)
    img(s, FIG / "p19_router_apply/budget.emf", (SW - iw) / 2, Inches(1.90), w=iw)
    # Colour key: the p04 channel schematics stand in for a plain swatch legend.
    # Ordered so each channel sits under the panel it owns — base/keep-air →
    # Airborne, droplet route → Droplet, close route → Close-contact.
    tw = Inches(1.35)
    for name, col, fn, gx in [("Airborne", GREEN, "airborne", Inches(2.36)),
                              ("Droplet", RED, "droplet", Inches(5.51)),
                              ("Close-contact", BLUE, "contact", Inches(8.68))]:
        img(s, FIG / f"p04_channels/{fn}.emf", gx, Inches(5.60), w=tw)
        txt(s, gx + Inches(1.47), Inches(6.02), Inches(1.5), Inches(0.32), name,
            size=13, bold=True, color=col)
        rect(s, gx + Inches(1.47), Inches(6.40), Inches(1.2), Inches(0.06), fill=col)


def p18b_router_result():
    s = slide()
    chrome(s, "③ Channel router — per-function SR", "3 — Improvements")
    subtitle(s, "Lifts exactly what it routes — net +0.8 pt")
    datanote(s, "change-ablation / BBOB-24 / dim 2 / n = 20 / "
                "functions whose SR@1e-10 moved", Inches(2.05))
    iw = Inches(9.8)
    img(s, FIG / "p22_router_result/sr.emf", (SW - iw) / 2, Inches(2.05), w=iw)


def p_best2_mech():
    s = slide()
    chrome(s, "④ Route-gated best2", "3 — Improvements")
    subtitle(s, "A 2nd difference vector — only on droplet-routed runs")
    GRAYA = RGBColor(0x6B, 0x74, 0x7E)
    # ── the condition, drawn as a branch (was a separate WHEN strip) ─────────
    cx = SW / 2
    cw = Inches(3.2)
    cy = Inches(1.84)                            # clears the subtitle at 1.32"
    ch = Inches(0.46)
    condition_box(s, cx - cw / 2, cy, cw, ch)
    paras(s, cx - cw / 2, cy, cw, ch, [
        {"align": PP_ALIGN.CENTER,
         "runs": [("DROPLET Route?", 13, True, DARK)]},
    ], anchor=MSO_ANCHOR.MIDDLE)
    # neutral fork → NO (base, grey, left) / YES (best2, red, right)
    lcx_in, rcx_in = 3.35, 9.93       # rcx: widest eq box lands on the margin
    lcx, rcx = Inches(lcx_in), Inches(rcx_in)
    eh_in, ey_in = 0.44, 3.40
    branch_fork(s, cx, cy + ch, [
        (lcx, Inches(1.5), "NO", GRAYA),
        (rcx, Inches(1.5), "YES", RED),
    ], Inches(ey_in))
    # each child= equation above its own panel — bold-italic mathtext (\bm),
    # the same symbols the panel below labels its individuals with
    def _eq(bx_in, ew, line_col, name):
        softbox(s, Inches(bx_in) - ew / 2, Inches(ey_in), ew, Inches(eh_in),
                fill=BG, line=line_col, line_w=1.3)
        _eq_center(s, name, bx_in, ey_in + eh_in / 2, folder="p23_best2")
    _eq(lcx_in, Inches(4.6), GRAYA, "eq_single")    # ≈ equation width + padding
    _eq(rcx_in, Inches(5.7), RED, "eq_best2")
    # the two mechanism panels, one per branch, each with the grey caption line
    # the other schematic pages carry (bottom clears the footer at 7.02")
    pw, capw = Inches(3.2), Inches(4.1)
    img(s, FIG / "p23_best2/single.emf", lcx - pw / 2, Inches(3.92), w=pw)
    img(s, FIG / "p23_best2/best2.emf", rcx - pw / 2, Inches(3.92), w=pw)
    for bx, cap in [(lcx, "off route — bit-identical to base"),
                    (rcx, "on route — 2nd donor pair escapes")]:
        txt(s, bx - capw / 2, Inches(6.72), capw, Inches(0.34), cap,
            size=14, color=GRAY, align=PP_ALIGN.CENTER)


def p_best2_conv():
    conv_result_slide(
        "④ Route-gated best2 — result",
        "Sharp ridges no longer trap runs",
        "best2_conv.npz", "p24_best2_result",
        [("a", "F13-SharpRidge"), ("b", "F14-DiffPowers")])


def p_best2_bar():
    s = slide()
    chrome(s, "④ Route-gated best2 — per-function SR", "3 — Improvements")
    subtitle(s, "F13/F14 jump to 100% — net +2.1 pt")
    datanote(s, "change-ablation / BBOB-24 / dim 2 / n = 20 / "
                "functions whose SR@1e-10 moved", Inches(2.05))
    iw = Inches(9.8)
    img(s, FIG / "p25_best2_bar/sr.emf", (SW - iw) / 2, Inches(2.05), w=iw)


def p_ladder():
    s = slide()
    chrome(s, "The four changes together", "3 — Improvements")
    subtitle(s, "SR@1e-10 by function, base → MC-ESO")
    datanote(s, "change-ablation / BBOB-24 / dim 2 / n = 20 / "
                "base → +restart → +floor → +router → +best2", Inches(2.1))
    iw = Inches(12.0)
    img(s, FIG / "p26_ladder/sr.emf", (SW - iw) / 2, Inches(2.1), w=iw)


def _family_conv_slide(title, subtitle_text, folder, rows):
    """Convergence-vs-existing-methods slide, p21 row layout: per function a
    small 2-D map + 3-D landscape (label above) and a large convergence panel
    on the right with the mean curve + semi-transparent ±1σ band per method."""
    s = slide()
    chrome(s, title, "4 — Comparison")
    # no subtitle — freed space becomes vertical margin around the two rows
    xm, wm = Inches(1.4), Inches(1.9)      # 2-D map (small square)
    xd, wd = Inches(3.45), Inches(2.0)     # 3-D landscape (small)
    xc, wc = Inches(5.6), Inches(6.35)     # convergence (large, wide)
    for (tag, name), y in zip(rows, [Inches(1.62), Inches(4.55)]):
        txt(s, xm, y, Inches(4.05), Inches(0.46), name, size=12, bold=True,
            color=DARK, align=PP_ALIGN.CENTER)
        img(s, FIG / f"{folder}/{tag}_map.emf", xm, y + Inches(0.4), w=wm)
        img(s, FIG / f"{folder}/{tag}_surf.emf", xd, y + Inches(0.46), w=wd)
        img(s, FIG / f"{folder}/{tag}_conv.png", xc, y + Inches(0.1), w=wc)


def p_family_conv1():
    _family_conv_slide(
        "Where baselines struggle — multimodal & ill-conditioned",
        "CMA-ES stalls on F03; PSO & SaVOA never reach on F11",
        "p27_family_conv",
        [("a", "F03 / Rastrigin (separable) — multimodal"),
         ("b", "F11 / Discus — ill-conditioned")])


def p_family_conv2():
    _family_conv_slide(
        "Where baselines struggle — bent valley & hardest case",
        "DE is slowest on F08; only MC-ESO solves F15",
        "p28_family_conv",
        [("c", "F08 / Rosenbrock — bent valley"),
         ("d", "F15 / Rastrigin (rotated) — hardest")])


# The five comparison methods, coloured to match the convergence curves in
# fig_family_conv (_CMP_COL). Shown once up front so the reader knows the field.
# One numbered reference list for the whole deck; the [n] markers on the method
# slides all resolve here. Rendered on p_references() just before the appendix.
_REFERENCES = [
    'A. Auger, N. Hansen. "A restart CMA-ES with increasing population size." IEEE CEC, 2005.',
    'J. Kennedy, R. Eberhart. "Particle swarm optimization." IEEE ICNN, 1995.',
    'R. Storn, K. Price. "Differential evolution." J. Global Optimization, 1997.',
    'Y.-C. Liang, J. R. Cuevas Juárez. "A self-adaptive virus optimization algorithm." Soft Computing, 2020.',
    'J. A. Nelder, R. Mead. "A simplex method for function minimization." The Computer Journal, 1965.',
    'N. Hansen, A. Ostermeier. "Completely derandomized self-adaptation in evolution strategies." Evolutionary Computation, 2001.',
    'N. Hansen. "Benchmarking a BI-population CMA-ES on the BBOB-2009 testbed." GECCO Workshop, 2009.',
    'R. Tanabe, A. Fukunaga. "Improving SHADE using linear population size reduction." IEEE CEC, 2014.',
    'B.-Y. Qu, P. N. Suganthan, J.-J. Liang. "DE with neighborhood mutation for multimodal optimization." IEEE TEVC, 2012.',
    'N. Hansen, A. Auger, R. Ros, O. Mersmann, T. Tušar, D. Brockhoff. "COCO: comparing continuous optimizers in a black-box setting." Optim. Methods & Software, 2021.',
]


def _tintc(hx, f):                              # blend hex toward white (f = colour amount)
    r, g, b = int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)
    return RGBColor(*(int(255 - (255 - v) * f) for v in (r, g, b)))


def _darkc(hx, f=0.72):                         # blend hex toward black
    r, g, b = int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)
    return RGBColor(*(int(v * f) for v in (r, g, b)))


# five comparison methods: name / ref# / colour / one-line what / struggle-role tag
_BASELINE_CARDS = [
    ("IPOP-CMA-ES", 1, "1F77B4", "CMA-ES with restarts + growing population",
     "weak", "multimodal"),
    ("PSO", 2, "FF7F0E", "Inertia-weighted particle swarm",
     "weak", "ill-conditioned"),
    ("DE", 3, "9467BD", "DE/rand/1/bin — the droplet channel's ancestor",
     "weak", "multimodal"),
    ("SaVOA", 4, "2CA02C", "Self-adaptive virus optimization",
     "weak", "ill-conditioned"),
    ("NM-Restart", 5, "8C564B", "Restarted Nelder–Mead simplex",
     "strong", "low-dim"),
]


def p_baselines_intro():
    s = slide()
    chrome(s, "The comparison set", "4 — Comparison")
    subtitle(s, "Five baselines — colours match the curves that follow")
    y, box_h, step = Inches(2.02), Inches(0.78), Inches(0.92)
    for name, ref, hexc, desc, kind, tag in _BASELINE_CARDS:
        col = _hexc(hexc)
        # row card with a colour accent bar on its left edge
        softbox(s, MARGIN, y, Inches(12.23), box_h, fill=WHITE, line=BORDER)
        rect(s, MARGIN + Inches(0.16), y + Inches(0.17), Inches(0.11),
             box_h - Inches(0.34), fill=col)
        paras(s, MARGIN + Inches(0.46), y, Inches(2.5), box_h,
              [{"runs": [(name, 15, True, col), (f"  [{ref}]", 11, False, GRAY)]}],
              anchor=MSO_ANCHOR.MIDDLE)
        txt(s, MARGIN + Inches(3.05), y, Inches(6.3), box_h, desc,
            size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        # right-side pill: where it struggles (or, for NM, where it's strong)
        pill_w, pill_h = Inches(2.55), Inches(0.44)
        px = MARGIN + Inches(12.23) - pill_w - Inches(0.18)
        py = y + (box_h - pill_h) / 2
        softbox(s, px, py, pill_w, pill_h, fill=_tintc(hexc, 0.16))
        txt(s, px, py, pill_w, pill_h, f"{kind}:  {tag}", size=11.5,
            bold=True, color=_darkc(hexc), align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(step))
    txt(s, MARGIN, Inches(6.74), Inches(12.2), Inches(0.3),
        "Others (CMA-ES / BIPOP / L-SHADE / NCDE)  →  appendix",
        size=10.5, italic=True, color=GRAY)


def p_references():
    s = slide()
    chrome(s, "References", "References")
    y = Inches(1.78)
    step = Inches(0.492)
    for i, cite in enumerate(_REFERENCES, start=1):
        paras(s, MARGIN, y, Inches(12.23), step,
              [{"runs": [(f"[{i}]", 12, True, RED_DK),
                         (f"   {cite}", 12, False, DARK)]}],
              anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(step))


def p_appendix_baselines():
    s = slide()
    chrome(s, "Appendix — the other baselines", "Appendix")
    subtitle(s, "In the full comparison; omitted from charts for legibility")
    rows = [
        ("CMA-ES", 6, "Plain covariance-matrix ES with multistart — the IPOP base "
         "without population doubling."),
        ("BIPOP-CMA-ES", 7, "CMA-ES restarts alternating large / small population "
         "regimes to keep the budgets balanced."),
        ("L-SHADE", 8, "Success-history adaptive DE with linear population "
         "reduction (CEC2014 winner)."),
        ("NCDE", 9, "Neighborhood-based crowding DE — a niching specialist, used "
         "only for the multi-optima (PR / MMOsr) comparison."),
    ]
    y = Inches(2.05)
    for name, ref, desc in rows:
        paras(s, MARGIN, y, Inches(2.7), Inches(0.7),
              [{"runs": [(name, 15, True, RED_DK), (f"  [{ref}]", 11, False, GRAY)]}],
              anchor=MSO_ANCHOR.MIDDLE)
        rect(s, Inches(3.25), y + Inches(0.04), Inches(0.02), Inches(0.62),
             fill=RULEC)
        txt(s, Inches(3.5), y, Inches(9.3), Inches(0.7), desc, size=13,
            color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(Inches(1.05)))


# ── appendix: the full 35-function landscape-signal table (router diagnostic) ─
# Measured on unmodified base MC-ESO (scripts/measure_channel_signals.py),
# n_runs=8, dim2, median over exploration generations. Source: docs/history.md.
# Rows in BBOB category order; highlighted cells = the router's firing
# thresholds (cond ≥ 3.0 → droplet, algA > 0.965 ∧ mgap > 0.36 → close).
_SIG_HEADER = ("function", "route", "cond", "PR", "algA", "offd",
               "divs", "kurt", "mgap", "nelt", "spil")
_SIG_GROUPS = [
    ("Separable — F01–05", [
        ("F01-Sphere",            "-",        "0.38", "1.71", "0.917", "0.268", "0.002", "−0.13", "0.301", "1.60", "4.0"),
        ("F02-EllipsoidalSep",    "droplet?", "5.62", "1.00", "1.000", "0.406", "0.005", "1.24",  "0.380", "1.95", "3.4"),
        ("F03-RastriginSep",      "-",        "1.03", "1.19", "0.994", "0.294", "0.018", "0.08",  "0.397", "2.50", "4.1"),
        ("F04-BucheRastrigin",    "close",    "0.91", "1.25", "0.988", "0.315", "0.010", "0.48",  "0.411", "2.03", "7.5"),
        ("F05-LinearSlope",       "-",        "2.28", "1.01", "0.999", "0.437", "0.003", "−0.09", "0.290", "1.71", "4.2"),
    ]),
    ("Moderate — F06–09", [
        ("F06-AttractiveSector",  "close?",   "1.26", "1.11", "0.957", "0.698", "0.002", "0.43",  "0.349", "1.48", "2.9"),
        ("F07-StepEllipsoidal",   "-",        "1.57", "1.05", "0.781", "0.944", "0.009", "−0.62", "0.241", "2.12", "5.0"),
        ("F08-Rosenbrock",        "-",        "2.56", "1.01", "0.897", "0.986", "0.007", "−0.18", "0.272", "2.12", "3.4"),
        ("F09-RosenbrockRot",     "-",        "2.55", "1.01", "0.962", "0.963", "0.006", "−0.07", "0.293", "2.00", "3.4"),
    ]),
    ("Ill-conditioned — F10–14", [
        ("F10-EllipsoidalRot",    "droplet?", "5.53", "1.00", "0.977", "1.000", "0.009", "−0.28", "0.243", "1.96", "3.1"),
        ("F11-Discus",            "droplet",  "5.42", "1.00", "0.953", "1.000", "0.009", "−0.26", "0.239", "1.99", "3.2"),
        ("F12-BentCigar",         "droplet",  "5.64", "1.00", "1.000", "0.956", "0.009", "−0.29", "0.271", "2.09", "3.8"),
        ("F13-SharpRidge",        "droplet",  "4.89", "1.00", "0.956", "1.000", "0.005", "0.13",  "0.257", "1.67", "6.0"),
        ("F14-DiffPowers",        "droplet",  "3.34", "1.00", "0.965", "0.993", "0.002", "−0.34", "0.246", "1.50", "3.6"),
    ]),
    ("Multimodal — F15–19", [
        ("F15-RastriginRot",      "keep-air", "1.07", "1.18", "0.861", "0.774", "0.018", "0.01",  "0.343", "2.33", "4.4"),
        ("F16-Weierstrass",       "close",    "1.63", "1.05", "0.994", "0.486", "0.050", "0.21",  "0.427", "3.35", "4.1"),
        ("F17-SchafferF7",        "keep-air", "0.95", "1.23", "0.974", "0.474", "0.003", "0.02",  "0.290", "1.72", "2.8"),
        ("F18-SchafferF7ill",     "?",        "2.63", "1.00", "0.974", "0.978", "0.008", "−0.24", "0.268", "1.95", "2.6"),
        ("F19-GriewankRosenbrock","keep-air", "0.73", "1.38", "0.951", "0.513", "0.101", "−0.14", "0.311", "4.29", "5.2"),
    ]),
    ("Weak structure — F20–24", [
        ("F20-Schwefel",          "keep-air", "1.16", "1.14", "0.956", "0.691", "0.024", "−0.27", "0.313", "2.38", "6.5"),
        ("F21-Gallagher101",      "-",        "1.31", "1.10", "0.799", "0.897", "0.017", "−0.48", "0.275", "2.56", "4.0"),
        ("F22-Gallagher21",       "-",        "2.74", "1.00", "0.900", "0.993", "0.012", "−0.56", "0.239", "2.29", "4.1"),
        ("F23-Katsuura",          "-",        "0.67", "1.43", "0.912", "0.435", "0.144", "−0.09", "0.282", "5.49", "3.4"),
        ("F24-LunacekRastrigin",  "keep-air", "0.64", "1.47", "0.876", "0.498", "0.076", "0.21",  "0.325", "4.31", "13.1"),
    ]),
    ("Custom (multi-optimum) — C01–11", [
        ("C01-Himmelblau",        "-",        "0.83", "1.30", "0.933", "0.547", "0.001", "1.32",  "0.603", "1.81", "2.4"),
        ("C02-SixHumpCamel",      "-",        "1.28", "1.21", "0.980", "0.609", "0.053", "−0.56", "0.574", "2.19", "3.6"),
        ("C03-Shubert",           "-",        "0.65", "1.45", "0.940", "0.385", "0.038", "1.15",  "0.533", "3.00", "3.5"),
        ("C04-FiveWell",          "-",        "0.46", "1.62", "0.915", "0.309", "0.002", "−0.13", "0.310", "1.99", "3.5"),
        ("C05-Eggholder",         "keep-air", "1.33", "1.12", "0.823", "0.754", "0.001", "0.75",  "0.394", "1.96", "11.6"),
        ("C06-Michalewicz",       "-",        "0.62", "1.46", "0.971", "0.255", "0.002", "0.44",  "0.363", "2.15", "3.9"),
        ("C07-BukinN6",           "-",        "4.15", "1.00", "0.993", "0.992", "0.001", "1.26",  "0.335", "1.67", "13.0"),
        ("C08-StyblinskiTang",    "-",        "0.49", "1.59", "0.935", "0.313", "0.001", "0.79",  "0.415", "1.65", "3.9"),
        ("C09-Easom",             "-",        "0.32", "1.78", "0.873", "0.252", "0.250", "−0.73", "0.252", "5.30", "7.6"),
        ("C10-SchafferN2",        "-",        "0.34", "1.76", "0.894", "0.258", "0.031", "−0.50", "0.248", "2.72", "4.0"),
        ("C11-DeJongF5",          "keep-air", "0.54", "1.53", "0.951", "0.300", "0.006", "−0.07", "0.340", "2.68", "7.8"),
    ]),
]
_SIG_LEGEND = [
    ("cond", "log10 eigenvalue ratio"), ("PR", "participation ratio [1,2]"),
    ("algA", "axis alignment [.71,1]"), ("offd", "corr. off-diag RMS [0,1]"),
    ("divs", "variance/span"), ("kurt", "excess kurtosis"),
    ("mgap", "max normalized gap"), ("nelt", "mean niche count"),
    ("spil", "spillover count"),
]
_SIG_CATBAND = RGBColor(0xE3, 0xE7, 0xEC)


def _sig_float(v):
    return float(v.replace("−", "-"))


def _route_color(route):
    if route.startswith("droplet"):
        return RED
    if route.startswith("close"):
        return BLUE
    if route == "keep-air":
        return GRN_DK
    return None                                  # "-" / "?": neutral


def _sig_cell_style(ci, val):
    """(fill, text-colour, bold) for a router-trigger cell, else None."""
    if ci == 2 and _sig_float(val) >= 3.0:                     # cond → droplet
        return _tintc("C0392B", 0.15), RED_DK, True
    if ci == 4 and _sig_float(val) > 0.965:                    # algA → close ①
        return _tintc("2E6DA4", 0.15), BLUE, True
    if ci == 8 and _sig_float(val) > 0.36:                     # mgap → close ②
        return _tintc("2E6DA4", 0.15), BLUE, True
    return None


def _signal_half_table(s, lx, t, groups, col_w, row_h, fs):
    total_w = Emu(sum(int(w) for w in col_w))
    ins = Inches(0.04)
    hdr_col = {2: RED_DK, 4: BLUE, 8: BLUE}     # the router's three signals
    y = t
    # header row
    rect(s, lx, y, total_w, row_h, fill=THEAD, line=BORDER, line_w=0.75)
    x = lx
    for ci, (name, cw) in enumerate(zip(_SIG_HEADER, col_w)):
        txt(s, x + ins, y, cw - 2 * ins, row_h, name, size=fs, bold=True,
            color=hdr_col.get(ci, GRAY),
            align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(int(x) + int(cw))
    y = Emu(int(y) + int(row_h))
    for gname, rows in groups:
        # category band
        rect(s, lx, y, total_w, row_h, fill=_SIG_CATBAND)
        txt(s, lx + Inches(0.08), y, total_w - Inches(0.16), row_h, gname,
            size=fs, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(row_h))
        for i, row in enumerate(rows):
            rect(s, lx, y, total_w, row_h, fill=TALT if i % 2 else WHITE)
            x = lx
            for ci, (cell, cw) in enumerate(zip(row, col_w)):
                col, bold = DARK, False
                st = _sig_cell_style(ci, cell) if ci >= 2 else None
                if st:
                    fill, col, bold = st
                    rect(s, x, y, cw, row_h, fill=fill)
                elif ci == 0:
                    bold = True
                elif ci == 1:
                    rc = _route_color(cell)
                    col, bold = (rc, True) if rc else (GRAY, False)
                txt(s, x + ins, y, cw - 2 * ins, row_h, cell, size=fs,
                    bold=bold, color=col,
                    align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
                x = Emu(int(x) + int(cw))
            y = Emu(int(y) + int(row_h))
    hline(s, lx, y, total_w, color=BORDER, weight=1.0)


def p_appendix_signals():
    s = slide()
    chrome(s, "Appendix — landscape-signal measurements", "Appendix")
    subtitle(s, "All 35 functions in BBOB category order — the diagnostic behind "
                "the channel router")
    # legend line 1: signal abbreviations
    leg = []
    for i, (k, v) in enumerate(_SIG_LEGEND):
        leg.append((("   " if i else "") + k, 9, True, DARK))
        leg.append((f" = {v}" + ("," if i < len(_SIG_LEGEND) - 1 else ""),
                    9, False, GRAY))
    leg += [("      values", 9, True, DARK),
            (" = median over exploration generations (σ > drilling)",
             9, False, GRAY)]
    # legend line 2: what the highlighted cells mean (the router's triggers)
    leg2 = [("highlighted cells", 9, True, DARK),
            (" = router triggers:   ", 9, False, GRAY),
            ("cond ≥ 3.0 → droplet", 9, True, RED_DK),
            ("  (checked first)   ", 9, False, GRAY),
            ("algA > 0.965 ∧ mgap > 0.36 → close", 9, True, BLUE),
            ("   else → ", 9, False, GRAY),
            ("keep-air (= base)", 9, True, GRN_DK)]
    paras(s, MARGIN, Inches(1.62), Inches(12.23), Inches(0.52),
          [{"runs": leg, "spacing": 1.15},
           {"runs": leg2, "space_before": 3}])
    # two side-by-side half-tables: BBOB F01–F19 (left) / F20–24 + Custom (right)
    col_w = [Inches(1.60), Inches(0.58)] + [Inches(0.425)] * 9
    row_h, tab_t, fs = Inches(0.198), Inches(2.24), 8
    datanote(s, "base MC-ESO (unmodified) / dim 2 / n = 8", tab_t)
    _signal_half_table(s, MARGIN, tab_t, _SIG_GROUPS[:4], col_w, row_h, fs)
    _signal_half_table(s, Inches(6.78), tab_t, _SIG_GROUPS[4:], col_w, row_h, fs)


# ══════════════════════════════════════════════════════════════════════════
# Assemble
# ══════════════════════════════════════════════════════════════════════════
title_slide()                                   # 1
divider(1, "Recap", "Where the 5/18 report left off")
p2_motivation()
p3_method_grid()                                 # channels ⊕ strategies, 3x2
p4_prev_result()
divider(2, "Direction", "Where to grow — and what we chose")
p7_direction()                                  # (diagnosis slide retired)
divider(3, "Improvements", "Four changes — what, why, and the numbers")
p8_timeline()
p9_waterfall()
# each improvement opens with the method map, highlighting the panel it refines
p3_method_grid("Restart  (spillover)", "① Informed restart", "3 — Improvements")
p12_restart()                                    # ① informed restart — real data
p12b_restart_math()                              # ① how it works (mechanism)
p13b_restart_why()                               # ① result — convergence (2 seeds)
p13c_restart_bar()                               # ① per-function SR bar
p3_method_grid("Close-contact", "② Adaptive anisotropy floor", "3 — Improvements")
p15_floor()                                      # ② adaptive floor — real data
p15c_floor_math()                                # ② the math (annotated equations)
p15b_floor_why()                                 # ② result — convergence (2 seeds)
p17b_floor_bar()                                 # ② per-function SR bar
p3_method_grid({"Droplet", "Close-contact", "Airborne"},
               "③ Channel router", "3 — Improvements")
p18_router()                                     # ③ route: shape / condition / budget (merged)
p18c_router_conv()                               # ③ result — convergence (2 seeds)
p18b_router_result()                             # ③ per-function SR bar
p3_method_grid("Droplet", "④ Route-gated best2", "3 — Improvements")
p_best2_mech()                                   # ④ mechanism (2nd difference)
p_best2_conv()                                   # ④ result — convergence (2 seeds)
p_best2_bar()                                    # ④ per-function SR bar
p_ladder()                                       # summary — cumulative SR ladder
p16_evals()                                      # summary — evals-to-success ladder (moved after ladder)
divider(4, "Comparison", "MC-ESO vs existing methods")
p_baselines_intro()                              # who the five comparison methods are
p_family_conv1()                                 # conv vs existing methods — F03 / F11
p_family_conv2()                                 # conv vs existing methods — F08 / F15
p10_methods()                                    # standing vs the unified set (table + Wilcoxon)
divider(5, "Multimodality", "The deferred axis — partial progress")
p18_multimodal_diag()                            # facade + scope/remaining (merged)
p3_method_grid("Strain coexistence", "Sequential niching", "5 — Multimodality")
p19_pr_vs_sr()                                   # niching — schematic + PR/SR evidence (merged)
divider(6, "Wrap-up", "Summary and what remains")
p21_summary()
p22_future()
p_references()                                    # numbered reference list [1]–[10]
p_appendix_baselines()                           # the other baselines (CMA/BIPOP/L-SHADE/NCDE)
p_appendix_signals()                             # full 35-function signal table (router diagnostic)

TOTAL_PAGES = len(prs.slides)

# replace the {TP} placeholder in page-number textboxes
for sl in prs.slides:
    for shp in sl.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                if "{TP}" in run.text:
                    run.text = run.text.replace("{TP}", str(TOTAL_PAGES))

out = OUTDIR / "20260714.pptx"
prs.save(str(out))
print("saved", out, "with", TOTAL_PAGES, "slides")
